"""Build the Task 3 presentation (.pptx) from Task3_Presentation_Plan_VI.md content.

Generates ``reports/Task3_Presentation.pptx`` — 33 slides (16:9), matching the
storyboard: title, agenda, anchor, the five Task-3 parts, the EDA block, and the
conclusion. Slide bodies carry the (A)/Làm-gì-Vì-sao-Kết-quả bullets and key
numbers; the (B) lời thoại goes into each slide's speaker notes; real figures from
``reports/figures`` are embedded where the plan specifies.

    python tools/build_task3_slides.py
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

try:
    from PIL import Image

    def _img_size(path: Path) -> tuple[int, int]:
        with Image.open(path) as im:
            return im.size
except Exception:  # pragma: no cover - PIL optional

    def _img_size(path: Path) -> tuple[int, int]:
        return (1600, 900)


ROOT = Path(__file__).resolve().parents[1]
FIG = ROOT / "reports" / "figures"
OUT = ROOT / "reports" / "Task3_Presentation.pptx"

# ── Palette ───────────────────────────────────────────────────────────────
NAVY = RGBColor(0x0F, 0x2A, 0x4F)
ACCENT = RGBColor(0xE8, 0x72, 0x2B)  # FPT-ish orange
GREEN = RGBColor(0x1E, 0x8E, 0x4E)
BLUE = RGBColor(0x1F, 0x6F, 0xB2)
AMBER = RGBColor(0xB9, 0x7A, 0x0C)
DARK = RGBColor(0x26, 0x2B, 0x33)
GREY = RGBColor(0x5B, 0x63, 0x70)
LIGHT = RGBColor(0xF2, 0xF4, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Calibri"
SW, SH = Inches(13.333), Inches(7.5)

LABELS = {"do": ("Làm gì", GREEN), "why": ("Vì sao", BLUE), "out": ("Kết quả", AMBER)}

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]


def _set(run, size, color, bold=False, italic=False, font=FONT):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font


def _box(slide, left, top, width, height):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    return tb, tf


def _rect(slide, left, top, width, height, fill, line=None):
    from pptx.enum.shapes import MSO_SHAPE

    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


def _footer(slide, n, presenter=""):
    _rect(slide, 0, 0, SW, Inches(0.16), NAVY)  # top accent bar
    _rect(slide, 0, 0, Inches(3.2), Inches(0.16), ACCENT)
    _, tf = _box(slide, Inches(0.55), Inches(7.06), Inches(9.0), Inches(0.34))
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "DSP391m · Nhóm 1 · Task 3 — Thu thập & Tiền xử lý dữ liệu"
    _set(r, 9, GREY)
    if presenter:
        r2 = p.add_run()
        r2.text = f"   ·   {presenter}"
        _set(r2, 9, ACCENT, bold=True)
    _, tf2 = _box(slide, Inches(12.2), Inches(7.06), Inches(0.9), Inches(0.34))
    tf2.paragraphs[0].alignment = PP_ALIGN.RIGHT
    rr = tf2.paragraphs[0].add_run()
    rr.text = str(n)
    _set(rr, 10, NAVY, bold=True)


def _notes(slide, text):
    if not text:
        return
    slide.notes_slide.notes_text_frame.text = text


def add_title(title, subtitle, meta, presenters):
    s = prs.slides.add_slide(BLANK)
    _rect(s, 0, 0, SW, SH, NAVY)
    _rect(s, 0, Inches(2.55), SW, Inches(0.07), ACCENT)
    _, tf = _box(s, Inches(0.9), Inches(2.7), Inches(11.5), Inches(2.2))
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    _set(r, 33, WHITE, bold=True)
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = subtitle
    _set(r2, 16, RGBColor(0xCF, 0xDC, 0xEC), italic=True)
    p2.space_before = Pt(10)
    _, tk = _box(s, Inches(0.9), Inches(1.5), Inches(11.0), Inches(0.7))
    rk = tk.paragraphs[0].add_run()
    rk.text = "BÁO CÁO TÁC VỤ DỮ LIỆU"
    _set(rk, 15, ACCENT, bold=True)
    _, tm = _box(s, Inches(0.9), Inches(5.3), Inches(11.5), Inches(1.7))
    for i, line in enumerate(meta):
        pp = tm.paragraphs[0] if i == 0 else tm.add_paragraph()
        rr = pp.add_run()
        rr.text = line
        _set(rr, 13, RGBColor(0xD7, 0xE0, 0xEC))
        pp.space_after = Pt(2)
    pp = tm.add_paragraph()
    rr = pp.add_run()
    rr.text = presenters
    _set(rr, 12.5, WHITE, bold=True)
    pp.space_before = Pt(6)


def add_divider(part, title, subtitle, presenter, n):
    s = prs.slides.add_slide(BLANK)
    _rect(s, 0, 0, SW, SH, NAVY)
    _rect(s, Inches(0.9), Inches(2.7), Inches(1.6), Inches(0.09), ACCENT)
    _, tk = _box(s, Inches(0.9), Inches(2.0), Inches(11), Inches(0.7))
    rk = tk.paragraphs[0].add_run()
    rk.text = part
    _set(rk, 15, ACCENT, bold=True)
    _, tf = _box(s, Inches(0.9), Inches(2.95), Inches(11.3), Inches(2.0))
    r = tf.paragraphs[0].add_run()
    r.text = title
    _set(r, 30, WHITE, bold=True)
    if subtitle:
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = subtitle
        _set(r2, 15, RGBColor(0xCF, 0xDC, 0xEC), italic=True)
        p2.space_before = Pt(8)
    _, tp = _box(s, Inches(0.9), Inches(5.6), Inches(11), Inches(0.6))
    rp = tp.paragraphs[0].add_run()
    rp.text = f"Trình bày: {presenter}"
    _set(rp, 14, WHITE, bold=True)
    # Divider = transition slide (intentionally unnumbered; content slides keep 1–33).


def _add_heading(s, kicker, title):
    _, tk = _box(s, Inches(0.55), Inches(0.32), Inches(12), Inches(0.4))
    rk = tk.paragraphs[0].add_run()
    rk.text = kicker
    _set(rk, 12, ACCENT, bold=True)
    _, tt = _box(s, Inches(0.55), Inches(0.66), Inches(12.2), Inches(0.95))
    rt = tt.paragraphs[0].add_run()
    rt.text = title
    _set(rt, 25, NAVY, bold=True)
    _rect(
        s,
        Inches(0.55),
        Inches(1.62),
        Inches(12.2),
        Inches(0.025),
        RGBColor(0xD8, 0xDF, 0xE8),
    )


def _render_bullets(tf, items):
    first = True
    for it in items:
        kind = it[0]
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(7)
        p.line_spacing = 1.04
        if kind in LABELS:
            label, col = LABELS[kind]
            dot = p.add_run()
            dot.text = "● "
            _set(dot, 13, col, bold=True)
            lab = p.add_run()
            lab.text = f"{label}: "
            _set(lab, 14.5, col, bold=True)
            txt = p.add_run()
            txt.text = it[1]
            _set(txt, 14.5, DARK)
        elif kind == "p":
            dot = p.add_run()
            dot.text = "▪ "
            _set(dot, 12, NAVY, bold=True)
            txt = p.add_run()
            txt.text = it[1]
            _set(txt, 14.5, DARK)
        elif kind == "s":
            p.level = 1
            dot = p.add_run()
            dot.text = "– "
            _set(dot, 12, GREY)
            txt = p.add_run()
            txt.text = it[1]
            _set(txt, 12.5, GREY)
            p.space_after = Pt(4)
        elif kind == "key":
            bar = p.add_run()
            bar.text = "▶ Số liệu chốt:  "
            _set(bar, 13.5, ACCENT, bold=True)
            txt = p.add_run()
            txt.text = it[1]
            _set(txt, 13.5, NAVY, bold=True)
            p.space_before = Pt(6)
        elif kind == "note":
            txt = p.add_run()
            txt.text = it[1]
            _set(txt, 12, GREY, italic=True)
            p.space_before = Pt(4)


def _add_image(s, path: Path, left, top, max_w, max_h, caption=""):
    w, h = _img_size(path)
    ratio = min(max_w / w, max_h / h)
    iw, ih = int(w * ratio), int(h * ratio)
    pic_left = left + (max_w - iw) // 2
    s.shapes.add_picture(str(path), pic_left, top, width=Emu(iw), height=Emu(ih))
    if caption:
        _, tf = _box(s, left, top + Emu(ih) + Inches(0.05), Emu(max_w), Inches(0.4))
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        r = tf.paragraphs[0].add_run()
        r.text = caption
        _set(r, 9.5, GREY, italic=True)


def _add_table(s, headers, rows, left, top, width, col_w=None, fontsize=11):
    nrows, ncols = len(rows) + 1, len(headers)
    height = Inches(0.36 * nrows)
    gtbl = s.shapes.add_table(nrows, ncols, left, top, width, height)
    tbl = gtbl.table
    if col_w:
        for i, cw in enumerate(col_w):
            tbl.columns[i].width = cw
    for j, htext in enumerate(headers):
        c = tbl.cell(0, j)
        c.fill.solid()
        c.fill.fore_color.rgb = NAVY
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        c.margin_top = c.margin_bottom = Pt(2)
        para = c.text_frame.paragraphs[0]
        r = para.add_run()
        r.text = htext
        _set(r, fontsize, WHITE, bold=True)
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            c = tbl.cell(i, j)
            c.fill.solid()
            c.fill.fore_color.rgb = WHITE if i % 2 else LIGHT
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            c.margin_top = c.margin_bottom = Pt(1)
            para = c.text_frame.paragraphs[0]
            r = para.add_run()
            r.text = str(val)
            _set(r, fontsize, NAVY if j == 0 else DARK, bold=(j == 0))
    return gtbl


def add_content(
    n,
    kicker,
    title,
    items,
    notes="",
    presenter="",
    image=None,
    caption="",
    table=None,
    text_w=12.2,
):
    s = prs.slides.add_slide(BLANK)
    _footer(s, n, presenter)
    _add_heading(s, kicker, title)
    if image:
        text_w = 6.85
    _, tf = _box(s, Inches(0.6), Inches(1.85), Inches(text_w), Inches(5.0))
    tf.vertical_anchor = MSO_ANCHOR.TOP
    _render_bullets(tf, items)
    if image:
        _add_image(
            s, image, Inches(7.7), Inches(1.95), Inches(5.15), Inches(4.6), caption
        )
    if table:
        _add_table(
            s,
            table["headers"],
            table["rows"],
            Inches(table.get("left", 0.6)),
            Inches(table["top"]),
            Inches(table.get("width", 9.5)),
            col_w=table.get("col_w"),
            fontsize=table.get("fs", 11),
        )
    _notes(s, notes)
    return s


# ════════════════════════════════════════════════════════════════════════
# SLIDES
# ════════════════════════════════════════════════════════════════════════

# S1 — Title
add_title(
    "Time-Aware Explainable ML cho phát hiện sớm sinh viên nguy cơ (OULAD)",
    "Thu thập & Tiền xử lý dữ liệu — Báo cáo Task 3",
    [
        "DSP391m · Đồ án Khoa học Dữ liệu · Nhóm 1 · Đại học FPT",
        "GVHD: Nguyễn Thị Hoàng Yến",
    ],
    "Thành viên: Sơn · Khoa · An · Đức · Phúc · Bình",
)

# S2 — Agenda
add_content(
    2,
    "MỤC LỤC",
    "Nội dung trình bày",
    [
        ("p", "Phần 1 — Xác định dữ liệu cần có  (Sơn)"),
        ("p", "Phần 2 — Phương pháp thu thập dữ liệu  (Phúc)"),
        ("p", "Phần 3 — Làm sạch dữ liệu  (Bình)"),
        ("p", "Khối EDA — Phân tích khám phá dữ liệu  (Bình)"),
        ("p", "Phần 4 — Chuẩn hoá & biến đổi dữ liệu  (Đức)"),
        ("p", "Phần 5 — Tách train/test & chống rò rỉ  (Khoa)"),
        ("note", "Mỗi mục trình bày theo: Làm gì – Vì sao – Kết quả đầu ra."),
    ],
    notes="Chào thầy/cô và cả lớp. Bài báo cáo Task 3 của nhóm gồm 5 phần theo đúng yêu cầu, "
    "cộng thêm một khối Phân tích khám phá dữ liệu (EDA). Mỗi mục chúng em trình bày theo "
    "công thức: làm gì, vì sao làm, và kết quả đầu ra là gì.",
)

# S3 — Anchor / big picture
add_content(
    3,
    "BỨC TRANH TỔNG THỂ",
    "Pipeline dữ liệu của nhóm",
    [
        (
            "do",
            "7 bảng OULAD thô → gộp → master_raw (32.593 × 33) → cắt theo 6 mốc thời gian → tiền xử lý chống rò rỉ → train/test.",
        ),
        ("why", "Slide 'neo' để luôn biết đang ở bước nào của quy trình."),
        (
            "out",
            "6 bộ dữ liệu dataset_t10…t100 + tập test 20% cố định; 19/19 kiểm thử rò rỉ ĐẠT.",
        ),
    ],
    notes="Đây là bức tranh tổng thể. Chúng em bắt đầu từ 7 bảng OULAD thô, gộp thành một bảng "
    "master 32.593 dòng 33 cột, cắt dữ liệu theo 6 mốc tiến độ khóa học, tiền xử lý theo "
    "trình tự chống rò rỉ, rồi tách train/test. Đầu ra là 6 bộ dữ liệu theo mốc thời gian và "
    "một tập kiểm tra cố định; toàn bộ được 19 kiểm thử tự động xác nhận không rò rỉ.",
    image=(FIG / "preprocessing_sequence.png"),
    caption="Trình tự pipeline tiền xử lý (preprocessing_sequence.png)",
)

# ── PHẦN 1 ───────────────────────────────────────────────────────────────
add_divider(
    "PHẦN 1",
    "Xác định dữ liệu cần có",
    "Mục tiêu · Biến mục tiêu · Nhóm dữ liệu · Đơn vị & phạm vi",
    "Sơn",
    4,
)

add_content(
    4,
    "PHẦN 1 · XÁC ĐỊNH DỮ LIỆU",
    "Bài toán & mục tiêu đề tài",
    [
        (
            "do",
            "Phát hiện SỚM sinh viên nguy cơ học kém trên VLE, và GIẢI THÍCH được mỗi dự đoán.",
        ),
        (
            "why",
            "Ba khoảng trống của nghiên cứu trước: hộp đen (opacity), dự đoán muộn (lateness), mất cân bằng lớp.",
        ),
        (
            "out",
            "Dữ liệu phải phục vụ RQ1 (dự đoán sớm), RQ2 (ổn định giải thích), RQ3 (mất cân bằng).",
        ),
    ],
    presenter="Sơn",
    notes="Đề tài giải quyết bài toán phát hiện sớm sinh viên có nguy cơ học kém, và quan trọng là "
    "phải giải thích được. Nghiên cứu trước có ba điểm yếu: mô hình hộp đen, dự đoán quá muộn, "
    "và mất cân bằng lớp. Vì vậy dữ liệu của Task 3 phải đủ để trả lời ba câu hỏi nghiên cứu.",
)

add_content(
    5,
    "PHẦN 1 · XÁC ĐỊNH DỮ LIỆU",
    "Biến mục tiêu (at_risk)",
    [
        ("do", "Phân loại nhị phân: at_risk suy ra từ final_result."),
        (
            "p",
            "at-risk (1) = Fail + Withdrawn  ·  not-at-risk (0) = Pass + Distinction.",
        ),
        ("why", "Gom theo 'có cần can thiệp hay không' — đúng mục tiêu cảnh báo sớm."),
        (
            "key",
            "at-risk 52,8% (17.208) vs not-at-risk 47,2% (15.385) — số THẬT, không phải 68/32.",
        ),
    ],
    presenter="Sơn",
    notes="Đây là bài toán phân loại nhị phân. Nhãn at_risk gộp Fail và Withdrawn; not-at-risk gồm "
    "Pass và Distinction — gom theo tiêu chí có cần can thiệp hay không. Tỉ lệ at-risk thật là "
    "52,8%, tức mất cân bằng nhẹ, không phải con số minh họa 68/32.",
    image=(FIG / "target_distribution.png"),
    caption="Phân phối biến mục tiêu (target_distribution.png)",
)

add_content(
    6,
    "PHẦN 1 · XÁC ĐỊNH DỮ LIỆU",
    "Ba nhóm dữ liệu đầu vào & vai trò",
    [
        ("do", "Ba nhóm đặc trưng:"),
        (
            "s",
            "Nhân khẩu học — gender, region, education, imd_band, age_band, disability…",
        ),
        ("s", "Tương tác/VLE — clickstream: tổng/loại click, số ngày hoạt động…"),
        ("s", "Kết quả đánh giá — điểm bài nộp, số bài nộp, cờ chưa nộp."),
        (
            "why",
            "Nhân khẩu = bối cảnh & công bằng; tương tác = tín hiệu hành vi sớm; kết quả = bằng chứng năng lực.",
        ),
        (
            "out",
            "Bộ đặc trưng chia 3 nhóm rõ ràng → thuận tiện diễn giải SHAP/LIME theo nhóm.",
        ),
    ],
    presenter="Sơn",
    notes="Dữ liệu đầu vào chia ba nhóm: nhân khẩu học cho bối cảnh và phân tích công bằng; tương "
    "tác clickstream là tín hiệu hành vi sớm; kết quả đánh giá là bằng chứng năng lực. Chia "
    "nhóm rõ giúp giải thích mô hình theo từng nhóm về sau.",
)

add_content(
    7,
    "PHẦN 1 · XÁC ĐỊNH DỮ LIỆU",
    "Đơn vị quan sát, phạm vi & dữ liệu loại trừ",
    [
        (
            "do",
            "Đơn vị quan sát: 1 dòng = 1 sinh viên–môn–kỳ (code_module × code_presentation × id_student).",
        ),
        (
            "do",
            "Phạm vi: 32.593 bản ghi · 22 môn–kỳ · 7 bảng · các kỳ 2013–2014 của Open University.",
        ),
        (
            "do",
            "Loại trừ: bỏ cột định danh; KHÔNG dùng date_unregistration làm đặc trưng; bỏ final_result gốc.",
        ),
        (
            "why",
            "Tách đơn vị quan sát & loại trừ sớm để tránh rò rỉ và nhân bản dữ liệu.",
        ),
    ],
    presenter="Sơn",
    notes="Mỗi dòng dữ liệu là một sinh viên trong một môn–kỳ. Phạm vi gồm hơn 32 nghìn bản ghi trên "
    "22 môn–kỳ. Chúng em loại sớm các cột định danh, biến date_unregistration (gây rò rỉ tương "
    "lai) và biến final_result gốc sau khi đã tạo nhãn.",
)

# ── PHẦN 2 ───────────────────────────────────────────────────────────────
add_divider(
    "PHẦN 2",
    "Phương pháp thu thập dữ liệu",
    "Nguồn · Cấu trúc 7 bảng · Liên kết · Ưu/Nhược/Rủi ro",
    "Phúc",
    8,
)

add_content(
    8,
    "PHẦN 2 · THU THẬP DỮ LIỆU",
    "Nguồn & loại dữ liệu",
    [
        (
            "do",
            "Nguồn: Open University Learning Analytics Dataset (OULAD), Kuzilek và cộng sự (2017).",
        ),
        (
            "do",
            "Loại: dữ liệu THỨ CẤP công khai, đã ẩn danh tại nguồn, giấy phép CC-BY 4.0.",
        ),
        (
            "why",
            "Đáp ứng đạo đức bằng trích dẫn đúng quy cách; không xử lý dữ liệu cá nhân nhạy cảm.",
        ),
        (
            "out",
            "Hợp pháp, tái lập, không cần thu thập sơ cấp — phù hợp môi trường học thuật.",
        ),
    ],
    presenter="Phúc",
    notes="Dữ liệu lấy từ bộ OULAD của Open University, công bố năm 2017. Đây là dữ liệu thứ cấp công "
    "khai, đã ẩn danh tại nguồn, cấp phép CC-BY 4.0, nên chỉ cần trích dẫn đúng là đáp ứng yêu "
    "cầu đạo đức.",
)

add_content(
    9,
    "PHẦN 2 · THU THẬP DỮ LIỆU",
    "Cách thu thập & cấu trúc 7 bảng",
    [
        (
            "do",
            "Tải 7 CSV → đặt vào data/raw/ → xác minh toàn vẹn bằng MD5 (data_manifest.txt).",
        ),
        (
            "do",
            "7 bảng: studentInfo, courses, studentRegistration, studentVle (10.655.280 dòng), vle, assessments, studentAssessment.",
        ),
        (
            "why",
            "Tách bảng nền khỏi bảng nặng → xử lý theo chunk (500k dòng/lần) tránh tràn bộ nhớ.",
        ),
        (
            "key",
            "Dữ liệu thô được 'đóng băng' + có checksum → mọi thành viên dùng đúng một bản.",
        ),
    ],
    presenter="Phúc",
    notes="Quy trình thu thập: tải 7 file CSV, đặt vào thư mục data/raw, rồi xác minh toàn vẹn bằng mã "
    "băm MD5. Bảng nặng nhất là studentVle với hơn 10,6 triệu dòng clickstream, nên chúng em đọc "
    "theo từng khối 500 nghìn dòng để tránh tràn bộ nhớ.",
)

add_content(
    10,
    "PHẦN 2 · THU THẬP DỮ LIỆU",
    "Mối liên hệ giữa các bảng",
    [
        ("do", "Khoá liên kết chính:"),
        (
            "s",
            "(code_module, code_presentation, id_student) — nối studentInfo ↔ registration ↔ VLE-agg ↔ performance.",
        ),
        ("s", "id_site — nối studentVle ↔ vle (lấy loại hoạt động)."),
        (
            "s",
            "id_assessment — nối studentAssessment ↔ assessments (lấy trọng số/hạn nộp).",
        ),
        (
            "why",
            "Hiểu khoá là điều kiện để gộp KHÔNG nhân bản (validate='many_to_one').",
        ),
        ("out", "Lược đồ quan hệ — nền tảng cho bước gộp ở Phần 3."),
    ],
    presenter="Phúc",
    notes="Các bảng liên kết qua ba loại khoá: khoá sinh viên–môn–kỳ nối các bảng chính; id_site nối "
    "clickstream với loại hoạt động; id_assessment nối bài nộp với trọng số và hạn nộp. Hiểu "
    "đúng khoá là điều kiện để gộp mà không nhân bản dữ liệu.",
)

add_content(
    11,
    "PHẦN 2 · THU THẬP DỮ LIỆU",
    "Đánh giá nguồn: phù hợp · ưu · hạn chế · rủi ro",
    [
        (
            "do",
            "Phù hợp: đủ 3 nhóm đặc trưng + dấu thời gian → cho phép dự đoán time-aware; 32k đủ lớn.",
        ),
        ("do", "Ưu điểm: chuẩn benchmark quốc tế, ẩn danh, tái lập, miễn phí."),
        (
            "why",
            "Hạn chế: thứ cấp 1 trường; click chỉ là số đếm (không nội dung); kỳ 2013–2014.",
        ),
        (
            "out",
            "Rủi ro: mất cân bằng nhẹ · thiếu cấu trúc · lệch phải mạnh · đa cộng tuyến · rò rỉ thời gian — đều có phương án ở Phần 3–5.",
        ),
    ],
    presenter="Phúc",
    notes="Nguồn phù hợp vì có đủ ba nhóm đặc trưng và có dấu thời gian để làm dự đoán theo mốc. Ưu "
    "điểm là chuẩn quốc tế, ẩn danh, tái lập, miễn phí. Hạn chế là dữ liệu một trường, click chỉ "
    "là số đếm. Các rủi ro như mất cân bằng, lệch, đa cộng tuyến, rò rỉ thời gian đều được xử lý "
    "ở các phần sau.",
)

# ── PHẦN 3 ───────────────────────────────────────────────────────────────
add_divider(
    "PHẦN 3",
    "Làm sạch dữ liệu",
    "Thiếu · Trùng · Kiểu · Ngoại lai · Gộp bảng",
    "Bình",
    12,
)

add_content(
    12,
    "PHẦN 3 · LÀM SẠCH",
    "Kiểm tra tổng quan dữ liệu thô",
    [
        ("do", "Sau gộp: 32.593 dòng × 33 cột; rà kiểu dữ liệu & cấu trúc."),
        ("why", "Bước 'khám tổng quát' trước khi can thiệp — biết quy mô & loại biến."),
        ("out", "Bảng tóm tắt số dòng – cột – kiểu, làm mốc cho các bước làm sạch."),
    ],
    presenter="Bình",
    notes="Sau khi gộp, bảng master có 32.593 dòng và 33 cột. Bước đầu tiên là khám tổng quát: số "
    "dòng, số cột, kiểu dữ liệu, để biết mình đang xử lý gì trước khi can thiệp.",
)

add_content(
    13,
    "PHẦN 3 · LÀM SẠCH",
    "Dữ liệu thiếu: phát hiện & phương án",
    [
        (
            "do",
            "Chỉ 3 cột thiếu; điểm/số bài thiếu do chưa nộp → điền 0 + cờ not_submitted.",
        ),
        (
            "why",
            "Mỗi loại thiếu (MCAR/MAR/MNAR) xử lý khác nhau; thiếu do 'chưa nộp' là TÍN HIỆU, không phải nhiễu.",
        ),
        ("out", "Sau xử lý: 0 giá trị khuyết ở mọi cột đặc trưng."),
    ],
    presenter="Bình",
    notes="Chỉ ba cột bị thiếu. date_unregistration thiếu 22.521 do đa số không rút môn — đây là thiếu "
    "cấu trúc, không dùng làm đặc trưng. imd_band thiếu 1.111 điền Unknown thành một nhóm riêng. "
    "date_registration thiếu 45 điền trung vị train. Điểm thiếu do chưa nộp điền 0 kèm cờ "
    "not_submitted — đây là tín hiệu chứ không phải nhiễu.",
    table={
        "headers": ["Cột", "Số thiếu", "Xử lý"],
        "rows": [
            ["date_unregistration", "22.521", "Bỏ (thiếu cấu trúc)"],
            ["imd_band", "1.111", "Điền 'Unknown'"],
            ["date_registration", "45", "Trung vị train"],
        ],
        "top": 4.55,
        "width": 6.9,
        "left": 0.6,
        "col_w": [Inches(3.0), Inches(1.5), Inches(2.4)],
        "fs": 11,
    },
    image=(FIG / "quality_missingness.png"),
    caption="quality_missingness.png",
)

add_content(
    14,
    "PHẦN 3 · LÀM SẠCH",
    "Trùng lặp & chuẩn hoá định dạng",
    [
        (
            "do",
            "Trùng: kiểm tra theo khoá sinh viên–môn–kỳ → drop_duplicates → 0 khoá trùng.",
        ),
        (
            "do",
            "Kiểu dữ liệu: danh mục numeric / ordinal / nominal / binary / indicator (VARIABLE_TYPES).",
        ),
        (
            "do",
            "Chuẩn hoá: strip text; đối chiếu từ điển (region 13, education 5, imd_band 10, age_band 3, gender/disability 2).",
        ),
        (
            "why",
            "Nhận diện đúng kiểu là điều kiện mã hoá đúng ở Phần 4; chuẩn hoá để 'Y ' ≠ 'Y' không xảy ra.",
        ),
        ("out", "Bộ dữ liệu nhất quán về định dạng & kiểu."),
    ],
    presenter="Bình",
    notes="Kiểm tra trùng theo khoá sinh viên–môn–kỳ: kết quả 0 khoá trùng. Lập danh mục kiểu biến để "
    "mã hoá đúng về sau. Chuẩn hoá định dạng bằng cách cắt khoảng trắng và đối chiếu số giá trị "
    "với từ điển dữ liệu.",
)

add_content(
    15,
    "PHẦN 3 · LÀM SẠCH",
    "Giá trị bất thường & ngoại lai: phát hiện",
    [
        (
            "do",
            "Dò bất thường logic (điểm ∈ [0–100]) + dò ngoại lai bằng quy tắc IQR (log_outliers).",
        ),
        (
            "why",
            "Clickstream lệch phải rất mạnh: clicks_resource skew ≈ 35, kurtosis ≈ 2.125; max_clicks_single_day max = 6.988.",
        ),
        (
            "out",
            "Bảng 7 biến nghi ngoại lai — bàn giao Bình → Đức (handoff_outliers_for_Duc.csv).",
        ),
    ],
    presenter="Bình",
    notes="Chúng em dò bất thường logic và dùng quy tắc IQR để phát hiện ngoại lai. Các biến click "
    "lệch phải rất mạnh, ví dụ clicks_resource có skew khoảng 35. Bảng các biến nghi ngoại lai "
    "được em bàn giao cho Đức để thống nhất chiến lược xử lý.",
    image=(FIG / "univariate_boxplots.png"),
    caption="univariate_boxplots.png",
)

add_content(
    16,
    "PHẦN 3 · LÀM SẠCH",
    "Xử lý bất thường & ngoại lai",
    [
        (
            "do",
            "Chiến lược theo biến: log1p cho 12 biến click lệch mạnh · winsorize (cắt 1% hai đầu) cho biến lệch vừa · none cho biến trong phạm vi tự nhiên.",
        ),
        (
            "why",
            "KHÔNG loại bỏ dòng nào: giá trị cực trị của sinh viên at-risk (vd học lại nhiều lần) là tín hiệu cần giữ.",
        ),
        ("out", "Mọi biến số 'thuần' hơn về phân phối mà vẫn giữ nguyên 32.593 dòng."),
    ],
    presenter="Bình",
    notes="Mỗi biến có chiến lược riêng: log1p cho biến click lệch mạnh, winsorize cắt 1% hai đầu cho "
    "biến lệch vừa, và giữ nguyên với biến trong phạm vi tự nhiên. Nguyên tắc quan trọng là không "
    "xoá dòng nào, vì giá trị cực trị của sinh viên nguy cơ là tín hiệu cần giữ.",
)

add_content(
    17,
    "PHẦN 3 · LÀM SẠCH",
    "Gộp bảng & bộ dữ liệu sạch cuối cùng",
    [
        (
            "do",
            "Gộp: studentInfo (nền) ← left join ← registration ← engagement ← performance ← courses.",
        ),
        (
            "why",
            "Nhật ký trước/sau chứng minh không nhân bản, không thất thoát (validate='many_to_one').",
        ),
        (
            "out",
            "master_raw = 32.593 × 33 — bộ dữ liệu sạch, một dòng một sinh viên–môn–kỳ.",
        ),
    ],
    presenter="Bình",
    notes="Bảy bảng được gộp bằng left join lên bảng nền studentInfo. Nhật ký số dòng trước và sau mỗi "
    "phép gộp đều giữ nguyên 32.593, chứng minh không nhân bản và không thất thoát. Đầu ra là "
    "bảng master sạch 32.593 dòng 33 cột.",
    table={
        "headers": ["Bước", "rows_before", "rows_after"],
        "rows": [
            ["studentInfo (nền)", "32.593", "32.593"],
            ["+ registration", "32.593", "32.593"],
            ["+ engagement", "32.593", "32.593"],
            ["+ performance", "32.593", "32.593"],
        ],
        "top": 4.6,
        "width": 9.0,
        "col_w": [Inches(4.2), Inches(2.4), Inches(2.4)],
        "fs": 12,
    },
)

# ── KHỐI EDA ─────────────────────────────────────────────────────────────
add_divider(
    "KHỐI EDA",
    "Phân tích khám phá dữ liệu",
    "Hiểu dữ liệu · Kiểm tra chất lượng · Tìm tín hiệu phân biệt",
    "Bình",
    18,
)

add_content(
    18,
    "EDA · MỤC TIÊU",
    "Mục tiêu EDA",
    [
        (
            "p",
            "Hiểu dữ liệu: 32.593 sinh viên, 33 biến (nhân khẩu học · tương tác · điểm số).",
        ),
        ("p", "Kiểm tra chất lượng: thiếu dữ liệu, kiểu dữ liệu, ngoại lai."),
        ("p", "Tìm tín hiệu phân biệt sinh viên at-risk vs not-at-risk."),
        ("p", "Phục vụ 2 quyết định: chọn đặc trưng & chống rò rỉ (leakage)."),
        ("key", "32.593 bản ghi · 33 cột · 3 nhóm biến."),
    ],
    presenter="Bình",
    notes="EDA của nhóm có 3 mục tiêu. Một là hiểu dữ liệu — hơn 32 nghìn bản ghi, 33 biến chia 3 "
    "nhóm: nhân khẩu học, tương tác (clickstream) và điểm số. Hai là kiểm tra chất lượng — "
    "thiếu, phân phối, ngoại lai. Ba, quan trọng nhất với đề tài phát hiện sớm, là tìm biến "
    "phân biệt được sinh viên nguy cơ. Mọi kết luận EDA phục vụ trực tiếp cho hai việc phía "
    "sau: chọn đặc trưng và đảm bảo không rò rỉ dữ liệu.",
)

add_content(
    19,
    "EDA · UNIVARIATE",
    "Phân tích phân phối (Univariate)",
    [
        ("p", "Công cụ: histogram + KDE, boxplot cho mỗi biến số."),
        ("p", "Phần lớn biến clickstream LỆCH PHẢI MẠNH (right-skewed)."),
        ("p", "Đuôi nặng → nhiều ngoại lai THẬT (sinh viên siêu tích cực), không xoá."),
        ("p", "Hệ quả: cần biến đổi log + chuẩn hoá; dùng kiểm định phi tham số."),
        (
            "key",
            "clicks_resource skew 34,71 · kurtosis 2.125,45 → lệch phải, đuôi nặng.",
        ),
    ],
    presenter="Bình",
    notes="Ở phân phối đơn biến, nhóm vẽ histogram kèm KDE và boxplot cho từng biến số. Rõ nhất là "
    "các biến đếm click lệch phải rất mạnh: clicks_resource có độ lệch skew tới 34,7 và độ nhọn "
    "kurtosis hơn 2.100 — đa số click ít, nhưng một nhóm nhỏ hoạt động cực nhiều tạo đuôi dài. "
    "Đây là ngoại lai thật, phản ánh hành vi học thật chứ không phải lỗi, nên nhóm không xoá mà "
    "xử lý bằng log1p và chuẩn hoá. Vì dữ liệu không chuẩn, các bước sau nhóm dùng kiểm định phi "
    "tham số Mann–Whitney thay cho t-test.",
    image=(FIG / "univariate_hist_kde.png"),
    caption="univariate_hist_kde.png",
)

add_content(
    20,
    "EDA · BIVARIATE",
    "Phân tích theo nhãn at-risk (Bivariate)",
    [
        ("p", "Phân phối lớp: at-risk 52,8% vs not-at-risk 47,2% → mất cân bằng nhẹ."),
        ("p", "So sánh 2 nhóm bằng Cohen's d + Mann–Whitney (hiệu chỉnh BH)."),
        (
            "p",
            "Top biến: days_since_last_activity d=2,55 · n_assessments_submitted 2,05 · weighted_score_to_date 1,96.",
        ),
        ("p", "19/19 biến số có ý nghĩa thống kê (q < 0,05)."),
        ("key", "days_since: 171 vs 14 ngày · n_assessments: 2,3 vs 8,6 bài."),
    ],
    presenter="Bình",
    notes="Phân phối nhãn: at-risk chiếm 52,8% — mất cân bằng nhẹ, là con số thật của bộ dữ liệu. Để "
    "tìm biến phân biệt tốt, với mỗi biến số nhóm so sánh hai nhóm bằng Cohen's d kèm Mann–"
    "Whitney, hiệu chỉnh đa kiểm định Benjamini–Hochberg. Ba biến mạnh nhất thuộc nhóm tương tác "
    "và điểm: số ngày kể từ lần hoạt động cuối (d khoảng 2,5 — rất lớn), số bài đã nộp, điểm có "
    "trọng số. Cụ thể, sinh viên at-risk trung bình 171 ngày không hoạt động so với chỉ 14 ngày "
    "ở nhóm còn lại; nộp trung bình 2,3 bài so với 8,6 bài. Toàn bộ 19/19 biến đều có ý nghĩa. "
    "Kết luận: tín hiệu nguy cơ nằm ở hành vi và điểm, không phải nhân khẩu học. Nếu được hỏi về "
    "nhân khẩu học: yếu, Cramér's V cao nhất chỉ khoảng 0,15.",
    image=(FIG / "bivariate_effect_sizes.png"),
    caption="bivariate_effect_sizes.png",
)

add_content(
    21,
    "EDA · MULTIVARIATE",
    "Tương quan & kiểm tra rò rỉ (Multivariate)",
    [
        ("p", "Heatmap Pearson (tuyến tính) + Spearman (đơn điệu, hợp dữ liệu lệch)."),
        (
            "p",
            "Tương quan với nhãn: days_since +0,78 · n_assessments −0,72 · weighted_score −0,71.",
        ),
        ("p", "Đa cộng tuyến: 2 cặp |r| ≥ 0,8 (n_days_active–total_clicks = 0,84)."),
        ("p", "Kiểm tra leakage: KHÔNG biến nào |r| ≥ 0,95 với nhãn → không rò rỉ."),
        ("key", "r(nhãn) max = 0,78 · leakage ≥0,95: 0 biến · Cramér's V ≤ 0,15."),
    ],
    presenter="Bình",
    notes="Bước đa biến dùng cả Pearson và Spearman — Spearman quan trọng vì dữ liệu lệch. Hai phát "
    "hiện chính. Một, tương quan với nhãn: days_since_last_activity +0,78, số bài nộp và điểm "
    "tương quan âm khoảng 0,71–0,72 — khớp đúng kết quả Cohen's d ở slide trước, hai phương pháp "
    "độc lập cho cùng kết luận. Hai, kiểm tra rò rỉ — bước bắt buộc: nhóm đặt ngưỡng nếu biến "
    "nào |r| ≥ 0,95 với nhãn thì nghi leakage; kết quả không biến nào vượt — mạnh nhất chỉ 0,78, "
    "hợp lý về giáo dục chứ không phải rò rỉ. Ngoài ra có 2 cặp đa cộng tuyến cao nên ở mô hình "
    "tuyến tính nhóm sẽ cân nhắc gộp/bỏ bớt. Biến phân loại liên hệ với nhãn đều yếu.",
    image=(FIG / "corr_with_target.png"),
    caption="corr_with_target.png",
)

add_content(
    22,
    "EDA · TIME-AWARE (RQ1)",
    "Phân tích theo thời gian & kết luận EDA",
    [
        ("p", "Đo Cohen's d từng biến tại 6 mốc 10%→100%; tín hiệu mạnh dần đều."),
        (
            "p",
            "Mốc đầu đạt d ≥ 0,8: n_days_active & days_since_last_activity từ 10% · điểm/số bài từ 20%.",
        ),
        (
            "p",
            "Withdrawn ngừng hoạt động sớm: median 233 ngày & 89 click vs 11 ngày & 1.425 click.",
        ),
        (
            "out",
            "Dự đoán sớm khả thi từ ~10–20% khóa học; ưu tiên đặc trưng tương tác + điểm.",
        ),
        (
            "key",
            "weighted_score d: 0,61 → 1,96 · n_assessments: 0,67 → 2,05 (10%→100%).",
        ),
    ],
    presenter="Bình",
    notes="Phần này gắn trực tiếp với RQ1: tín hiệu xuất hiện từ giai đoạn nào? Nhóm tính lại Cohen's d "
    "cho từng biến tại cả 6 mốc. Khả năng phân biệt tăng dần đều theo thời gian — ví dụ "
    "weighted_score_to_date tăng từ 0,61 ở mốc 10% lên gần 1,96 ở cuối khóa, và "
    "n_assessments_submitted từ 0,67 lên 2,05, đúng bằng giá trị ở slide bivariate. Quan trọng "
    "cho phát hiện sớm: ngay mốc 10%, n_days_active đã đạt mức phân biệt mạnh; đến mốc 20%, cả "
    "điểm và số bài nộp đều mạnh. Một minh chứng nữa cho cách xử lý nhãn Withdrawn: sinh viên rút "
    "môn có median 233 ngày không hoạt động và chỉ 89 click, so với 1.425 click ở nhóm an toàn. "
    "Tóm lại EDA cho ba kết luận: tín hiệu nằm ở tương tác và điểm; không có rò rỉ; và dự đoán "
    "sớm khả thi từ khoảng 20% tiến độ. LƯU Ý: tránh trích 'days_since 1,56 ở mốc 100%' cạnh "
    "'d=2,55' của slide trước — hai số tính trên hai bộ khác nhau (master vs dataset_t100).",
    image=(FIG / "time_discrimination_curve.png"),
    caption="time_discrimination_curve.png",
)

# ── PHẦN 4 ───────────────────────────────────────────────────────────────
add_divider(
    "PHẦN 4",
    "Chuẩn hoá & biến đổi dữ liệu",
    "Mã hoá · log1p · Chuẩn hoá · Đặc trưng · Mất cân bằng",
    "Đức",
    23,
)

add_content(
    23,
    "PHẦN 4 · CHUẨN HOÁ & BIẾN ĐỔI",
    "Mã hoá biến phân loại",
    [
        ("do", "4 chiến lược theo bản chất biến (xem bảng)."),
        ("why", "Mã hoá sai (one-hot cho biến thứ bậc) sẽ XOÁ thông tin thứ tự."),
        ("out", "Biến chữ → số, mô hình & XAI đọc được."),
    ],
    presenter="Đức",
    notes="Mã hoá phân loại theo bốn chiến lược tuỳ bản chất biến: thứ bậc dùng OrdinalEncoder để giữ "
    "thứ tự; danh định dùng OneHotEncoder; nhị phân dùng BinaryEncoder; biến chỉ báo giữ nguyên. "
    "Nếu dùng sai, ví dụ one-hot cho biến thứ bậc, sẽ làm mất thông tin thứ tự.",
    table={
        "headers": ["Loại", "Biến", "Phương pháp"],
        "rows": [
            ["Thứ bậc", "highest_education, imd_band, age_band", "OrdinalEncoder"],
            ["Danh định", "region, code_module, code_presentation", "OneHotEncoder"],
            ["Nhị phân", "gender, disability", "BinaryEncoder (0/1)"],
            ["Chỉ báo", "not_submitted", "passthrough"],
        ],
        "top": 3.9,
        "width": 12.1,
        "col_w": [Inches(2.0), Inches(6.6), Inches(3.5)],
        "fs": 12,
    },
)

add_content(
    24,
    "PHẦN 4 · CHUẨN HOÁ & BIẾN ĐỔI",
    "Biến đổi phân phối lệch (log1p)",
    [
        ("do", "Áp log1p cho các biến click lệch phải mạnh (đã nêu ở Phần 3)."),
        (
            "why",
            "Giảm đuôi nặng → ổn định mô hình tuyến tính & khoảng cách; bằng chứng skew ~13–35 ở EDA.",
        ),
        ("out", "Đặc trưng click có phân phối cân đối hơn, không mất dòng."),
    ],
    presenter="Đức",
    notes="Với các biến click lệch phải mạnh, chúng em áp dụng biến đổi log1p để giảm đuôi nặng, giúp "
    "ổn định mô hình tuyến tính. Bằng chứng là độ lệch skew từ 13 đến 35 thấy ở phần EDA đơn biến.",
    image=(FIG / "univariate_hist_kde.png"),
    caption="Phân phối lệch phải → log1p",
)

add_content(
    25,
    "PHẦN 4 · CHUẨN HOÁ & BIẾN ĐỔI",
    "Chuẩn hoá thang đo (StandardScaler)",
    [
        ("do", "StandardScaler đưa biến số về trung bình 0, độ lệch chuẩn 1."),
        (
            "why",
            "Tổng click (0–hàng nghìn) áp đảo điểm số (0–100) → bắt buộc cho Logistic Regression & ANN.",
        ),
        (
            "out",
            "CHỈ fit trên train, transform cho cả train/test (in scaler.mean_ chỉ tính từ train).",
        ),
    ],
    presenter="Đức",
    notes="StandardScaler đưa biến số về trung bình 0 và độ lệch chuẩn 1. Cần thiết vì tổng click lên "
    "hàng nghìn áp đảo điểm số 0–100 về biên độ, đặc biệt với Logistic Regression và mạng nơ-ron. "
    "Điểm mấu chốt chống rò rỉ: chỉ fit trên train, rồi transform cho cả train và test.",
)

add_content(
    26,
    "PHẦN 4 · CHUẨN HOÁ & BIẾN ĐỔI",
    "Tạo đặc trưng mới & kiểm tra ý nghĩa",
    [
        (
            "do",
            "Tương tác: total_clicks, n_days_active, 8× clicks_<loại>, max/mean clicks, days_since_last_activity.",
        ),
        (
            "do",
            "Kết quả: mean_score_to_date, weighted_score_to_date, n_assessments_submitted, cờ not_submitted.",
        ),
        (
            "why",
            "Mỗi đặc trưng gắn giả thuyết: 'ngừng hoạt động lâu / bỏ nộp bài ⇒ nguy cơ cao'.",
        ),
        (
            "out",
            "Kiểm tra ý nghĩa bằng |Cohen's d|: days_since 2,55 · n_assessments 2,05 · weighted_score 1,96.",
        ),
    ],
    presenter="Đức",
    notes="Chúng em tạo các đặc trưng phái sinh từ clickstream và bài nộp, mỗi đặc trưng gắn một giả "
    "thuyết: ngừng hoạt động lâu hoặc bỏ nộp bài thì nguy cơ cao. Ý nghĩa được kiểm chứng bằng "
    "Cohen's d ở phần EDA — các đặc trưng mới này phân biệt hai lớp rất mạnh.",
    image=(FIG / "bivariate_effect_sizes.png"),
    caption="|Cohen's d| theo đặc trưng",
)

add_content(
    27,
    "PHẦN 4 · CHUẨN HOÁ & BIẾN ĐỔI",
    "Mất cân bằng dữ liệu & phương án",
    [
        (
            "do",
            "Kiểm tra: at-risk 52,8% / not-at-risk 47,2% → tỷ số mất cân bằng 1,12 (nhẹ).",
        ),
        (
            "why",
            "Bỏ sót sinh viên nguy cơ là sai lầm đắt nhất → chỉ số chính là PR-AUC & recall lớp at-risk.",
        ),
        (
            "out",
            "Phương án RQ3: so sánh no-resampling / class-weight / SMOTE / ADASYN; CHỈ tái lấy mẫu trên train.",
        ),
    ],
    presenter="Đức",
    notes="Dữ liệu mất cân bằng nhẹ với tỷ số 1,12. Dù nhẹ, bỏ sót sinh viên nguy cơ là sai lầm đắt "
    "nhất, nên chỉ số chính là PR-AUC và recall trên lớp at-risk chứ không phải accuracy. Câu hỏi "
    "RQ3 sẽ so sánh các chiến lược SMOTE, ADASYN, class-weight, và chỉ tái lấy mẫu trên tập train.",
)

# ── PHẦN 5 ───────────────────────────────────────────────────────────────
add_divider(
    "PHẦN 5",
    "Tách train/test & chống rò rỉ",
    "Tỷ lệ · Phân tầng + nhóm · Time-aware · Lưu & kiểm tra",
    "Khoa",
    28,
)

add_content(
    28,
    "PHẦN 5 · TÁCH & CHỐNG RÒ RỈ",
    "Mục đích, tỷ lệ & phương pháp chia",
    [
        ("do", "Tỷ lệ 80% train / 20% test (TEST_SIZE = 0.2)."),
        (
            "do",
            "Phương pháp: StratifiedGroupKFold — phân tầng theo at_risk + gom nhóm theo id_student.",
        ),
        (
            "why",
            "Một SV học nhiều môn–kỳ ⇒ gom theo id_student để KHÔNG cho cùng người ở cả train lẫn test (rò rỉ nhóm).",
        ),
        ("out", "Test ≈ 6.489 dòng / 5.756 sinh viên · train ≈ 26.104 dòng."),
    ],
    presenter="Khoa",
    notes="Chia 80% train, 20% test. Phương pháp StratifiedGroupKFold vừa phân tầng theo nhãn at_risk "
    "vừa gom nhóm theo id_student. Lý do gom theo sinh viên: một sinh viên học nhiều môn–kỳ, nếu "
    "chia theo dòng thì cùng một người sẽ xuất hiện ở cả train và test, gây rò rỉ nhóm.",
)

add_content(
    29,
    "PHẦN 5 · TÁCH & CHỐNG RÒ RỈ",
    "Kiểm soát rò rỉ: 2 trục",
    [
        (
            "do",
            "Trục thời gian: cut_at_checkpoint() chỉ giữ sự kiện ≤ ngày mốc; cutoff_day = round(length × t/100).",
        ),
        (
            "do",
            "Trục đặc trưng: mọi bộ học (encoder, scaler, resampling) CHỈ fit trên train.",
        ),
        (
            "why",
            "Mô phỏng đúng 'thông tin có tại thời điểm dự đoán' — không nhìn tương lai.",
        ),
        (
            "out",
            "6 bộ dataset_t10…t100, cùng danh sách sinh viên; nhãn cố định 52,8% qua các mốc.",
        ),
    ],
    presenter="Khoa",
    notes="Rò rỉ kiểm soát theo hai trục. Trục thời gian: hàm cut_at_checkpoint chỉ giữ sự kiện đến "
    "ngày mốc, với ba quy tắc loại bài nộp và click sau mốc. Trục đặc trưng: mọi bộ học chỉ fit "
    "trên train. Kết quả là 6 bộ dữ liệu theo mốc dùng chung một danh sách sinh viên.",
    image=(FIG / "time_discrimination_curve.png"),
    caption="Sức phân biệt tăng dần theo mốc (RQ1)",
)

add_content(
    30,
    "PHẦN 5 · TÁCH & CHỐNG RÒ RỈ",
    "Trình tự tiền xử lý đúng & áp dụng train→test",
    [
        (
            "do",
            "Trình tự: chia train/test → điền khuyết → ngoại lai → mã hoá/chuẩn hoá → tái lấy mẫu.",
        ),
        (
            "do",
            "Tham số (median, ngưỡng winsorize, mean/std, danh mục encoder) học từ train → transform lên test.",
        ),
        (
            "why",
            "Fit trên toàn bộ dữ liệu trước khi chia ⇒ test 'rò rỉ' vào train ⇒ kết quả ảo.",
        ),
        (
            "out",
            "Phân bố nhãn sau chia: train 0,53 / test 0,52 — chênh ≤ 0,02, đạt phân tầng.",
        ),
    ],
    presenter="Khoa",
    notes="Trình tự tiền xử lý bắt buộc: chia trước, rồi mới điền khuyết, xử lý ngoại lai, mã hoá, "
    "chuẩn hoá, và cuối cùng tái lấy mẫu. Mọi tham số học từ train rồi áp lên test. Nếu fit trên "
    "toàn bộ dữ liệu trước khi chia thì test rò rỉ vào train, cho kết quả ảo. Sau khi chia, tỉ lệ "
    "nhãn train và test chênh nhau dưới 0,02.",
)

add_content(
    31,
    "PHẦN 5 · TÁCH & CHỐNG RÒ RỈ",
    "Lưu & kiểm tra lại: rò rỉ đã bị chặn",
    [
        (
            "do",
            "Lưu: test_student_ids.csv (commit vào repo) + parquet train/test; split_report.csv.",
        ),
        (
            "do",
            "Kiểm tra lại: 0 sinh viên trùng giữa train/test; không lệch kiểu/cột; nhãn cân đối.",
        ),
        (
            "why",
            "Test cố định, dùng lại y hệt qua cả 6 mốc ⇒ 6 điểm hiệu năng SO SÁNH ĐƯỢC.",
        ),
        (
            "key",
            "tests/test_leakage.py: 19/19 kiểm thử ĐẠT — bằng chứng tự động không rò rỉ.",
        ),
    ],
    presenter="Khoa",
    notes="Định nghĩa tập test được lưu cố định và commit vào repo, kèm báo cáo split_report. Kiểm tra "
    "lại cho thấy 0 sinh viên trùng giữa train và test. Test cố định dùng lại qua cả 6 mốc nên 6 "
    "điểm hiệu năng so sánh được với nhau. Toàn bộ được 19 kiểm thử tự động xác nhận không rò rỉ.",
)

# ── KẾT LUẬN ─────────────────────────────────────────────────────────────
add_content(
    32,
    "KẾT LUẬN",
    "Tóm tắt Task 3",
    [
        (
            "do",
            "Năm mục Task 3 đều hoàn thành với sản phẩm & số liệu cụ thể (xem bảng).",
        ),
    ],
    presenter="An",
    notes="Tóm lại, năm mục của Task 3 đều hoàn thành với sản phẩm và số liệu cụ thể: xác định dữ liệu, "
    "thu thập, làm sạch, chuẩn hoá và biến đổi, và tách train/test.",
    table={
        "headers": ["Mục", "Kết quả của nhóm"],
        "rows": [
            ["1. Xác định dữ liệu", "at_risk nhị phân 52,8%; 3 nhóm; 32.593 sv–môn–kỳ"],
            ["2. Thu thập", "OULAD CC-BY 4.0; 7 bảng; 10,6 triệu click; có MD5"],
            [
                "3. Làm sạch",
                "0 trùng; log1p+winsorize, 0 dòng bị xoá; master 32.593×33",
            ],
            [
                "4. Chuẩn hoá & biến đổi",
                "Ordinal+OneHot+Binary; StandardScaler; d≤2,55; RQ3",
            ],
            ["5. Tách train/test", "80/20 group+stratified; 6 mốc; 19/19 test đạt"],
        ],
        "top": 2.7,
        "width": 12.2,
        "col_w": [Inches(3.3), Inches(8.9)],
        "fs": 12.5,
    },
)

s33 = add_content(
    33,
    "KẾT LUẬN",
    "Đầu ra & bước tiếp theo",
    [
        (
            "out",
            "Đầu ra Task 3: master_raw + 6 bộ dataset_t10…t100 + tập test cố định + pipeline tái lập (seed=42, MD5, Restart & Run All).",
        ),
        (
            "do",
            "Bước tiếp: benchmark mô hình theo từng mốc (RQ1) → xử lý mất cân bằng (RQ3) → SHAP/LIME & độ ổn định giải thích (RQ2).",
        ),
        ("note", "Cảm ơn thầy/cô và cả lớp — nhóm sẵn sàng trả lời câu hỏi."),
    ],
    presenter="An",
    notes="Đầu ra của Task 3 là bảng master, sáu bộ dữ liệu theo mốc, một tập test cố định, và một "
    "pipeline tái lập hoàn toàn. Bước tiếp theo là benchmark mô hình theo từng mốc cho RQ1, xử lý "
    "mất cân bằng cho RQ3, và thêm lớp giải thích SHAP/LIME cùng đo độ ổn định cho RQ2. Em xin "
    "cảm ơn và nhóm sẵn sàng trả lời câu hỏi.",
)
# closing emphasis line
_, tf = _box(s33, Inches(0.6), Inches(6.0), Inches(12.2), Inches(0.8))
r = tf.paragraphs[0].add_run()
r.text = "Cảm ơn đã lắng nghe!"
_set(r, 22, ACCENT, bold=True)


OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(OUT))
print(f"Saved {OUT}  ({len(prs.slides._sldIdLst)} slides)")
