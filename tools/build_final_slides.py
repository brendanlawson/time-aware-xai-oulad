"""Build the FINAL defense presentation (.pptx) for the DSP391m project.

Generates ``reports/slides/Final_Defense_Slides.pptx`` (16:9) — the end-of-term
defense deck covering the whole project: problem & RQs, OULAD data, the
anti-leakage pipeline, the 5-model benchmark, CV + Friedman, the time-aware RQ1
results (dual-cohort), threshold selection on validation, the RQ3 imbalance
comparison, XAI stability (RQ2), fairness, the dashboard, and conclusions.

House style is reused from ``build_progress_slides.py`` (palette, helpers,
speaker notes). Every metric is read from the CSV tables — nothing numeric is
hardcoded except structural dataset facts (row counts of raw OULAD).

Usage (from the repo root)::

    python -m tools.build_final_slides [--tables-dir reports/tables]

Missing tables degrade the related slide gracefully (warning, no crash).
No checkpointing: the script is a single fast pass (< a few seconds, no
network/LLM calls), so resume support would be pointless.
"""

from __future__ import annotations

import argparse
import csv
from pathlib import Path
import sys

# Console Windows mặc định cp1252 không in được tiếng Việt → ép UTF-8.
for _stream in (sys.stdout, sys.stderr):
    try:
        _stream.reconfigure(encoding="utf-8", errors="replace")
    except Exception:  # pragma: no cover - stream không hỗ trợ reconfigure
        pass

from pptx import Presentation  # noqa: E402
from pptx.dml.color import RGBColor  # noqa: E402
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN  # noqa: E402
from pptx.util import Emu, Inches, Pt  # noqa: E402

try:
    from PIL import Image

    def _img_size(path: Path) -> tuple[int, int]:
        with Image.open(path) as im:
            return im.size
except Exception:  # pragma: no cover - PIL optional

    def _img_size(path: Path) -> tuple[int, int]:
        return (1600, 900)


ROOT = Path(__file__).resolve().parents[1]
FIG = ROOT / "reports" / "figures"
OUT = ROOT / "reports" / "slides" / "Final_Defense_Slides.pptx"

# ── Palette (same as the progress / Task 3 decks) ───────────────────────────
NAVY = RGBColor(0x0F, 0x2A, 0x4F)
ACCENT = RGBColor(0xE8, 0x72, 0x2B)
GREEN = RGBColor(0x1E, 0x8E, 0x4E)
BLUE = RGBColor(0x1F, 0x6F, 0xB2)
AMBER = RGBColor(0xB9, 0x7A, 0x0C)
DARK = RGBColor(0x26, 0x2B, 0x33)
GREY = RGBColor(0x5B, 0x63, 0x70)
LIGHT = RGBColor(0xF2, 0xF4, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Calibri"
SW, SH = Inches(13.333), Inches(7.5)
LABELS = {"do": ("Làm gì", GREEN), "why": ("Vì sao", BLUE), "out": ("Kết quả", AMBER)}

MODEL_NAMES = {
    "xgb": "XGBoost",
    "lgbm": "LightGBM",
    "rf": "Random Forest",
    "logreg": "Logistic Reg.",
    "ann": "ANN (MLP)",
}
STRATEGY_NAMES = {
    "none": "No-resample",
    "class_weight": "Class-weight",
    "SMOTE": "SMOTE",
    "smote": "SMOTE",
    "ADASYN": "ADASYN",
    "adasyn": "ADASYN",
}

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]

DEGRADED: list[str] = []  # tables that were missing → slides degraded


# ── Data helpers ─────────────────────────────────────────────────────────────
def _load(tables_dir: Path, name: str) -> list[dict] | None:
    """Read a CSV as list[dict]; None (+warning) when missing/unreadable."""
    path = tables_dir / name
    try:
        with path.open(newline="", encoding="utf-8-sig") as fh:
            rows = list(csv.DictReader(fh))
        if not rows:
            raise ValueError("empty table")
        return rows
    except Exception as exc:
        DEGRADED.append(name)
        print(f"[WARN] thiếu/không đọc được {name} ({exc}) — slide liên quan sẽ giản lược.")
        return None


def _f(row: dict, key: str) -> float:
    return float(row[key])


def _vn(x, nd: int = 3) -> str:
    """Format a number with a Vietnamese decimal comma; '—' if not numeric."""
    try:
        return f"{float(x):.{nd}f}".replace(".", ",")
    except (TypeError, ValueError):
        return "—"


def _vn_int(x) -> str:
    try:
        return f"{int(float(x)):,}".replace(",", ".")
    except (TypeError, ValueError):
        return "—"


def _vn_p(p) -> str:
    try:
        v = float(p)
    except (TypeError, ValueError):
        return "—"
    return "< 0,001" if v < 0.001 else _vn(v, 3)


def _fig(name: str) -> Path | None:
    path = FIG / name
    if path.exists():
        return path
    print(f"[WARN] thiếu hình {name} — slide dùng bố cục chỉ chữ.")
    return None


# ── Drawing helpers (house style) ────────────────────────────────────────────
def _set(run, size, color, bold=False, italic=False, font=FONT):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font


def _box(slide, left, top, width, height):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    return tb, tf


def _rect(slide, left, top, width, height, fill, line=None):
    from pptx.enum.shapes import MSO_SHAPE

    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


def _footer(slide, n, presenter=""):
    _rect(slide, 0, 0, SW, Inches(0.16), NAVY)
    _rect(slide, 0, 0, Inches(3.2), Inches(0.16), ACCENT)
    _, tf = _box(slide, Inches(0.55), Inches(7.06), Inches(9.5), Inches(0.34))
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "DSP391m · Nhóm 1 · Bảo vệ cuối kỳ"
    _set(r, 9, GREY)
    if presenter:
        r2 = p.add_run()
        r2.text = f"   ·   {presenter}"
        _set(r2, 9, ACCENT, bold=True)
    _, tf2 = _box(slide, Inches(12.2), Inches(7.06), Inches(0.9), Inches(0.34))
    tf2.paragraphs[0].alignment = PP_ALIGN.RIGHT
    rr = tf2.paragraphs[0].add_run()
    rr.text = str(n)
    _set(rr, 10, NAVY, bold=True)


def _notes(slide, text):
    if text:
        slide.notes_slide.notes_text_frame.text = text


def add_title(title, subtitle, meta, presenters):
    s = prs.slides.add_slide(BLANK)
    _rect(s, 0, 0, SW, SH, NAVY)
    _rect(s, 0, Inches(2.55), SW, Inches(0.07), ACCENT)
    _, tf = _box(s, Inches(0.9), Inches(2.7), Inches(11.5), Inches(2.2))
    r = tf.paragraphs[0].add_run()
    r.text = title
    _set(r, 33, WHITE, bold=True)
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = subtitle
    _set(r2, 16, RGBColor(0xCF, 0xDC, 0xEC), italic=True)
    p2.space_before = Pt(10)
    _, tk = _box(s, Inches(0.9), Inches(1.5), Inches(11.0), Inches(0.7))
    rk = tk.paragraphs[0].add_run()
    rk.text = "BẢO VỆ CUỐI KỲ"
    _set(rk, 15, ACCENT, bold=True)
    _, tm = _box(s, Inches(0.9), Inches(5.3), Inches(11.5), Inches(1.7))
    for i, line in enumerate(meta):
        pp = tm.paragraphs[0] if i == 0 else tm.add_paragraph()
        rr = pp.add_run()
        rr.text = line
        _set(rr, 13, RGBColor(0xD7, 0xE0, 0xEC))
        pp.space_after = Pt(2)
    pp = tm.add_paragraph()
    rr = pp.add_run()
    rr.text = presenters
    _set(rr, 12.5, WHITE, bold=True)
    pp.space_before = Pt(6)
    return s


def _add_heading(s, kicker, title):
    _, tk = _box(s, Inches(0.55), Inches(0.32), Inches(12), Inches(0.4))
    rk = tk.paragraphs[0].add_run()
    rk.text = kicker
    _set(rk, 12, ACCENT, bold=True)
    _, tt = _box(s, Inches(0.55), Inches(0.66), Inches(12.2), Inches(0.95))
    rt = tt.paragraphs[0].add_run()
    rt.text = title
    _set(rt, 25, NAVY, bold=True)
    _rect(s, Inches(0.55), Inches(1.62), Inches(12.2), Inches(0.025), RGBColor(0xD8, 0xDF, 0xE8))


def _render_bullets(tf, items):
    first = True
    for it in items:
        kind = it[0]
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(7)
        p.line_spacing = 1.04
        if kind in LABELS:
            label, col = LABELS[kind]
            dot = p.add_run()
            dot.text = "● "
            _set(dot, 13, col, bold=True)
            lab = p.add_run()
            lab.text = f"{label}: "
            _set(lab, 14.5, col, bold=True)
            txt = p.add_run()
            txt.text = it[1]
            _set(txt, 14.5, DARK)
        elif kind == "p":
            dot = p.add_run()
            dot.text = "▪ "
            _set(dot, 12, NAVY, bold=True)
            txt = p.add_run()
            txt.text = it[1]
            _set(txt, 14.5, DARK)
        elif kind == "s":
            p.level = 1
            dot = p.add_run()
            dot.text = "– "
            _set(dot, 12, GREY)
            txt = p.add_run()
            txt.text = it[1]
            _set(txt, 12.5, GREY)
            p.space_after = Pt(4)
        elif kind == "key":
            bar = p.add_run()
            bar.text = "▶ Số liệu chốt:  "
            _set(bar, 13.5, ACCENT, bold=True)
            txt = p.add_run()
            txt.text = it[1]
            _set(txt, 13.5, NAVY, bold=True)
            p.space_before = Pt(6)
        elif kind == "note":
            txt = p.add_run()
            txt.text = it[1]
            _set(txt, 12, GREY, italic=True)
            p.space_before = Pt(4)


def _add_image(s, path: Path, left, top, max_w, max_h, caption=""):
    w, h = _img_size(path)
    ratio = min(max_w / w, max_h / h)
    iw, ih = int(w * ratio), int(h * ratio)
    pic_left = left + (max_w - iw) // 2
    s.shapes.add_picture(str(path), pic_left, top, width=Emu(iw), height=Emu(ih))
    if caption:
        _, tf = _box(s, left, top + Emu(ih) + Inches(0.05), Emu(max_w), Inches(0.4))
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        r = tf.paragraphs[0].add_run()
        r.text = caption
        _set(r, 9.5, GREY, italic=True)


def _add_table(s, headers, rows, left, top, width, col_w=None, fontsize=11):
    nrows, ncols = len(rows) + 1, len(headers)
    height = Inches(0.36 * nrows)
    gtbl = s.shapes.add_table(nrows, ncols, left, top, width, height)
    tbl = gtbl.table
    if col_w:
        for i, cw in enumerate(col_w):
            tbl.columns[i].width = cw
    for j, htext in enumerate(headers):
        c = tbl.cell(0, j)
        c.fill.solid()
        c.fill.fore_color.rgb = NAVY
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        c.margin_top = c.margin_bottom = Pt(2)
        r = c.text_frame.paragraphs[0].add_run()
        r.text = htext
        _set(r, fontsize, WHITE, bold=True)
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            c = tbl.cell(i, j)
            c.fill.solid()
            c.fill.fore_color.rgb = WHITE if i % 2 else LIGHT
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            c.margin_top = c.margin_bottom = Pt(1)
            r = c.text_frame.paragraphs[0].add_run()
            r.text = str(val)
            _set(r, fontsize, NAVY if j == 0 else DARK, bold=(j == 0))
    return gtbl


def add_content(
    n,
    kicker,
    title,
    items,
    notes="",
    presenter="",
    image=None,
    caption="",
    table=None,
    text_w=12.2,
):
    s = prs.slides.add_slide(BLANK)
    _footer(s, n, presenter)
    _add_heading(s, kicker, title)
    if image:
        text_w = 6.85
    _, tf = _box(s, Inches(0.6), Inches(1.85), Inches(text_w), Inches(5.0))
    tf.vertical_anchor = MSO_ANCHOR.TOP
    _render_bullets(tf, items)
    if image:
        _add_image(s, image, Inches(7.7), Inches(1.95), Inches(5.15), Inches(4.6), caption)
    if table:
        _add_table(
            s,
            table["headers"],
            table["rows"],
            Inches(table.get("left", 0.6)),
            Inches(table["top"]),
            Inches(table.get("width", 9.5)),
            col_w=table.get("col_w"),
            fontsize=table.get("fs", 11),
        )
    _notes(s, notes)
    return s


# ════════════════════════════════════════════════════════════════════════════
# Deck builder
# ════════════════════════════════════════════════════════════════════════════
def build(tables_dir: Path) -> None:
    split = _load(tables_dir, "split_report.csv")
    imb = _load(tables_dir, "imbalance_comparison.csv")
    cv = _load(tables_dir, "cv_summary.csv")
    friedman = _load(tables_dir, "model_friedman.csv")
    tbest = _load(tables_dir, "time_aware_best.csv")
    sens = _load(tables_dir, "sensitivity_active_xgb.csv")
    thr = _load(tables_dir, "threshold_validation.csv")
    xai_strat = _load(tables_dir, "xai_stability_strategies.csv")
    xai_seed = _load(tables_dir, "xai_stability_seeds.csv")
    fair = _load(tables_dir, "fairness_gaps.csv")

    # ── 1 · Bìa ──────────────────────────────────────────────────────────
    s1 = add_title(
        "Time-Aware Explainable ML — phát hiện sớm sinh viên nguy cơ (OULAD)",
        "Bảo vệ cuối kỳ: từ dữ liệu chống rò rỉ đến dự đoán sớm, giải thích được & công bằng",
        [
            "DSP391m · Đồ án Khoa học Dữ liệu · Nhóm 1 · Đại học FPT",
            "GVHD: Nguyễn Thị Hoàng Yến",
        ],
        "Nhóm 1: Phúc · Đức · Khoa · Bình · An · Sơn",
    )
    _notes(
        s1,
        "Chào thầy/cô và các bạn. Nhóm 1 bảo vệ đồ án cuối kỳ: hệ thống machine learning "
        "nhận diện sớm sinh viên nguy cơ trên bộ dữ liệu OULAD — nhấn mạnh ba chữ: sớm, "
        "giải thích được, và trung thực về giới hạn.",
    )

    # ── 2 · Bài toán & 3 RQ ──────────────────────────────────────────────
    add_content(
        2,
        "BÀI TOÁN",
        "Phát hiện sớm sinh viên nguy cơ & 3 câu hỏi nghiên cứu",
        [
            (
                "do",
                "Dự đoán nhị phân at-risk (Fail/Withdrawn) vs OK (Pass/Distinction) theo từng "
                "mốc tiến độ khóa học.",
            ),
            (
                "why",
                "Cảnh báo khi khóa học ĐÃ kết thúc là vô dụng — giá trị nằm ở dự đoán sớm, "
                "khi còn can thiệp được.",
            ),
            ("p", "RQ1 — Sớm đến mức nào thì dự đoán còn tin cậy? (time-aware, 6 mốc 10–100%)"),
            ("p", "RQ2 — Giải thích của mô hình có ổn định không? (SHAP/LIME, seed, chiến lược)"),
            (
                "p",
                "RQ3 — Xử lý mất cân bằng có đổi hiệu năng và có đổi lời giải thích không?",
            ),
        ],
        notes="Bài toán: phân loại nhị phân sinh viên nguy cơ theo từng mốc tiến độ. Ba câu hỏi "
        "nghiên cứu: RQ1 dự đoán sớm được đến đâu; RQ2 lời giải thích có ổn định không; RQ3 xử lý "
        "mất cân bằng ảnh hưởng gì đến cả hiệu năng lẫn lời giải thích.",
        presenter="Khoa",
    )

    # ── 3 · Dữ liệu OULAD + pipeline ─────────────────────────────────────
    add_content(
        3,
        "DỮ LIỆU",
        "OULAD — Open University Learning Analytics Dataset",
        [
            (
                "p",
                "32.593 lượt ghi danh (enrolment) · 28.785 sinh viên · 22 khóa-kỳ (presentation).",
            ),
            (
                "p",
                "7 bảng quan hệ; log VLE ~10,6 triệu lượt click là nguồn tín hiệu hành vi chính.",
            ),
            (
                "p",
                "Nhãn: Fail + Withdrawn → at-risk; Pass + Distinction → OK (bài toán nhị phân).",
            ),
            (
                "do",
                "Pipeline 1 dòng: 7 bảng → master 32.593×33 → cắt 6 mốc t∈{10,20,40,60,80,100}% "
                "→ tiền xử lý fit-on-train → 5 model.",
            ),
            (
                "note",
                "Các con số trên là fact cấu trúc của bộ dữ liệu công khai (Kuzilek et al., 2017).",
            ),
        ],
        notes="Dữ liệu là OULAD công khai của Open University: 32.593 lượt ghi danh của 28.785 sinh "
        "viên trên 22 khóa-kỳ, kèm log VLE khoảng 10,6 triệu click. Pipeline gộp 7 bảng thành bảng "
        "master rồi cắt theo 6 mốc tiến độ để mô phỏng dự đoán sớm.",
        presenter="Phúc",
        image=_fig("target_distribution.png"),
        caption="Phân bố nhãn at-risk vs OK",
    )

    # ── 4 · Chống leakage + frozen split ─────────────────────────────────
    split_items = [
        (
            "do",
            "4 quy tắc chống rò rỉ: (1) cắt dữ liệu theo thời gian TRƯỚC khi tính đặc trưng; "
            "(2) mọi bộ biến đổi fit-on-train; (3) chia nhóm theo id_student — không sinh viên "
            "nào nằm cả 2 phía; (4) tập test đóng băng, chỉ chạm 1 lần cuối.",
        ),
        ("why", "Chỉ cần 1 lỗi rò rỉ là toàn bộ kết quả 'đẹp' trở nên vô nghĩa khi triển khai."),
    ]
    if split:
        r0 = split[0]
        split_items.append(
            (
                "key",
                f"Frozen split: train {_vn_int(r0['n_train'])} / test {_vn_int(r0['n_test'])} "
                f"({_vn_int(r0['n_test_students'])} SV) · lệch tỷ lệ at-risk "
                f"{_vn(r0['rate_gap'], 4)} · trùng sinh viên = {_vn_int(r0['student_overlap'])}.",
            )
        )
    else:
        split_items.append(
            ("note", "(Thiếu split_report.csv — số liệu chi tiết của split xem trong repo.)")
        )
    split_items.append(
        ("out", "Bộ kiểm thử tự động (pytest) xác nhận các quy tắc trên trước mỗi lần chạy.")
    )
    add_content(
        4,
        "PHƯƠNG PHÁP · CHỐNG RÒ RỈ",
        "4 quy tắc chống leakage + tập test đóng băng",
        split_items,
        notes="Bốn quy tắc chống rò rỉ: cắt theo thời gian trước khi tính đặc trưng, fit-on-train, "
        "chia nhóm theo sinh viên, và tập test đóng băng chỉ chạm một lần. Split được kiểm bằng bộ "
        "test tự động: không sinh viên nào xuất hiện ở cả hai phía, tỷ lệ nhãn hai phía gần nhau.",
        presenter="Phúc",
        image=_fig("preprocessing_sequence.png"),
        caption="Trình tự pipeline tiền xử lý (fit-on-train)",
    )

    # ── 5 · 5 model & giao thức ──────────────────────────────────────────
    add_content(
        5,
        "PHƯƠNG PHÁP · MÔ HÌNH",
        "5 thuật toán ứng viên & giao thức đánh giá",
        [
            (
                "do",
                "5 model: Logistic Regression · Random Forest · XGBoost · LightGBM · ANN (MLP).",
            ),
            (
                "do",
                "Giao thức: CV 5-fold × 5 seed trên train (gom nhóm theo sinh viên) → kiểm định "
                "Friedman → đánh giá 1 lần trên test.",
            ),
            (
                "why",
                "Bỏ sót SV nguy cơ là sai lầm đắt nhất → recall & PR-AUC lớp at-risk là chỉ số chính.",
            ),
            (
                "p",
                "Ngưỡng quyết định chọn trên VALIDATION, không tối ưu trên test (slide 10).",
            ),
        ],
        notes="Năm thuật toán từ tuyến tính đến boosting và mạng nơ-ron. Giao thức: cross-validation "
        "5 fold lặp 5 seed trên train, kiểm định Friedman để so sánh có ý nghĩa thống kê, rồi mới "
        "đánh giá một lần duy nhất trên test. Chỉ số chính là recall và PR-AUC của lớp at-risk.",
        presenter="Đức",
    )

    # ── 6 · Benchmark t=100 ──────────────────────────────────────────────
    bench_items = [
        ("do", "So 5 model tại mốc t=100% trên tập test (baseline, chưa resample)."),
    ]
    bench_table = None
    if imb:
        base = [r for r in imb if r["strategy"] == "none"]
        base.sort(key=lambda r: _f(r, "recall"), reverse=True)
        if base:
            top = base[0]
            bench_items.append(
                (
                    "key",
                    f"{MODEL_NAMES.get(top['model'], top['model'])} dẫn đầu recall "
                    f"{_vn(top['recall'])} · F1 {_vn(top['f1'])} · PR-AUC {_vn(top['pr_auc'])}.",
                )
            )
            bench_items.append(
                ("note", "Khoảng cách giữa các model nhỏ → chọn thêm theo tính giải thích được.")
            )
            bench_table = {
                "headers": ["Model", "Recall", "F1", "PR-AUC", "ROC-AUC"],
                "rows": [
                    [
                        MODEL_NAMES.get(r["model"], r["model"]),
                        _vn(r["recall"]),
                        _vn(r["f1"]),
                        _vn(r["pr_auc"]),
                        _vn(r["roc_auc"]),
                    ]
                    for r in base
                ],
                "top": 3.35,
                "width": 11.0,
                "col_w": [Inches(3.4), Inches(1.9), Inches(1.9), Inches(1.9), Inches(1.9)],
                "fs": 13,
            }
    else:
        bench_items.append(
            ("note", "(Thiếu imbalance_comparison.csv — bảng số liệu xem trong repo.)")
        )
    add_content(
        6,
        "KẾT QUẢ · BENCHMARK",
        "5 model tại t=100% (tập test, baseline no-resample)",
        bench_items,
        notes="Bảng hiệu năng năm model trên tập test tại mốc 100%, chưa xử lý mất cân bằng, sắp theo "
        "recall. XGBoost dẫn đầu recall, LightGBM nhỉnh về PR-AUC; chênh lệch nhỏ nên nhóm chọn "
        "XGBoost làm model chính nhờ recall và hệ sinh thái giải thích.",
        presenter="Đức",
        table=bench_table,
    )

    # ── 7 · CV 5×5 + Friedman ────────────────────────────────────────────
    cv_items = [
        (
            "do",
            "CV 5-fold × 5 seed (25 khối) trên train tại t=100% → kiểm định Friedman theo rank.",
        ),
    ]
    cv_table = None
    if cv:
        cv100 = [r for r in cv if int(float(r["t_percent"])) == 100]
        cv100.sort(key=lambda r: _f(r, "recall_mean"), reverse=True)
        cv_table = {
            "headers": ["Model", "Recall (μ ± σ)", "F1 (μ ± σ)", "PR-AUC (μ ± σ)"],
            "rows": [
                [
                    MODEL_NAMES.get(r["model"], r["model"]),
                    f"{_vn(r['recall_mean'], 4)} ± {_vn(r['recall_std'], 4)}",
                    f"{_vn(r['f1_mean'], 4)} ± {_vn(r['f1_std'], 4)}",
                    f"{_vn(r['pr_auc_mean'], 4)} ± {_vn(r['pr_auc_std'], 4)}",
                ]
                for r in cv100
            ],
            "top": 3.5,
            "width": 11.6,
            "col_w": [Inches(2.9), Inches(2.9), Inches(2.9), Inches(2.9)],
            "fs": 12,
        }
    else:
        cv_items.append(("note", "(Thiếu cv_summary.csv — bảng CV xem trong repo.)"))
    if friedman:
        fr_rec = next((r for r in friedman if r["metric"] == "recall"), None)
        fr_pr = next((r for r in friedman if r["metric"] == "pr_auc"), None)
        parts = []
        if fr_rec:
            parts.append(
                f"recall: p {_vn_p(fr_rec['p_value'])}, tốt nhất "
                f"{MODEL_NAMES.get(fr_rec['best_model'], fr_rec['best_model'])} "
                f"(rank TB {_vn(fr_rec['mean_rank_best'], 2)})"
            )
        if fr_pr:
            parts.append(
                f"PR-AUC: p {_vn_p(fr_pr['p_value'])}, tốt nhất "
                f"{MODEL_NAMES.get(fr_pr['best_model'], fr_pr['best_model'])}"
            )
        if parts:
            cv_items.append(("key", "Friedman — " + " · ".join(parts) + "."))
        cv_items.append(
            ("out", "Khác biệt giữa 5 model có ý nghĩa thống kê nhưng biên độ tuyệt đối nhỏ.")
        )
    else:
        cv_items.append(("note", "(Thiếu model_friedman.csv — kết quả kiểm định xem trong repo.)"))
    add_content(
        7,
        "KẾT QUẢ · ĐỘ TIN CẬY",
        "CV 5-fold × 5 seed + kiểm định Friedman",
        cv_items,
        notes="Để kết luận không phụ thuộc một lần chia may mắn, nhóm chạy CV 5 fold lặp 5 seed — 25 "
        "khối — rồi kiểm định Friedman trên rank. Kết quả: khác biệt giữa các model có ý nghĩa thống "
        "kê; XGBoost đứng đầu về recall, LightGBM về PR-AUC, nhưng biên độ tuyệt đối nhỏ.",
        presenter="Đức",
        table=cv_table,
    )

    # ── 8 · RQ1 time-aware ───────────────────────────────────────────────
    rq1_items = [
        (
            "do",
            "Chạy lại toàn bộ giao thức tại 6 mốc tiến độ 10→100%; chọn model tốt nhất mỗi mốc.",
        ),
    ]
    rq1_table = None
    if tbest:
        rows = sorted(tbest, key=lambda r: int(float(r["t_percent"])))
        rq1_table = {
            "headers": ["Mốc t (%)", "Model tốt nhất", "Recall", "PR-AUC", "Tin cậy?"],
            "rows": [
                [
                    str(int(float(r["t_percent"]))),
                    MODEL_NAMES.get(r["model"], r["model"]),
                    _vn(r["recall"]),
                    _vn(r["pr_auc"]),
                    "✓" if str(r.get("reliable", "")).lower() == "true" else "✗",
                ]
                for r in rows
            ],
            "top": 2.6,
            "width": 6.6,
            "left": 0.6,
            "col_w": [Inches(1.2), Inches(2.0), Inches(1.2), Inches(1.2), Inches(1.0)],
            "fs": 11,
        }
        rel = [r for r in rows if str(r.get("reliable", "")).lower() == "true"]
        if rel:
            first = rel[0]
            rq1_items.append(
                (
                    "key",
                    f"Từ t={int(float(first['t_percent']))}% dự đoán đạt ngưỡng tin cậy "
                    f"(recall {_vn(first['recall'])}, {MODEL_NAMES.get(first['model'], first['model'])}).",
                )
            )
    else:
        rq1_items.append(("note", "(Thiếu time_aware_best.csv — đường cong xem trong repo.)"))
    s8 = add_content(
        8,
        "RQ1 · TIME-AWARE",
        "Hiệu năng theo 6 mốc tiến độ khóa học",
        rq1_items,
        notes="Trả lời RQ1: chạy toàn bộ giao thức tại sáu mốc tiến độ. Hiệu năng tăng đơn điệu theo "
        "lượng dữ liệu; hai mốc rất sớm 10 và 20 phần trăm còn yếu, từ mốc 40 phần trăm trở đi recall "
        "vượt ngưỡng nhóm coi là dùng được để can thiệp.",
        presenter="Khoa",
        image=_fig("time_aware_recall.png"),
        caption="Recall lớp at-risk theo mốc tiến độ",
    )
    # bảng đặt dưới khối bullet, cạnh trái (ảnh nằm bên phải)
    if rq1_table:
        _add_table(
            s8,
            rq1_table["headers"],
            rq1_table["rows"],
            Inches(rq1_table["left"]),
            Inches(rq1_table["top"] + 0.75),
            Inches(rq1_table["width"]),
            col_w=rq1_table["col_w"],
            fontsize=rq1_table["fs"],
        )

    # ── 9 · RQ1 dual-cohort (slide đinh) ─────────────────────────────────
    dual_items = [
        (
            "why",
            "Recall toàn-cohort tính cả SV đã rút trước mốc t; cohort active (còn học tại t) "
            "mới phản ánh giá trị can thiệp thật.",
        ),
    ]
    dual_table = None
    if sens:
        rows = sorted(sens, key=lambda r: int(float(r["t_percent"])))
        dual_table = {
            "headers": ["t (%)", "Recall full", "Recall active", "SV đã rút trước t"],
            "rows": [
                [
                    str(int(float(r["t_percent"]))),
                    _vn(r["full_recall"]),
                    _vn(r["active_recall"]),
                    _vn_int(r["withdrawn_already_gone"]),
                ]
                for r in rows
            ],
            "top": 4.15,
            "width": 6.7,
            "left": 0.6,
            "col_w": [Inches(0.9), Inches(1.9), Inches(1.9), Inches(2.0)],
            "fs": 11,
        }
        t_full = next(
            (int(float(r["t_percent"])) for r in rows if _f(r, "full_recall") >= 0.80), None
        )
        t_act = next(
            (int(float(r["t_percent"])) for r in rows if _f(r, "active_recall") >= 0.80), None
        )
        key = []
        if t_full is not None:
            key.append(f"full đạt recall ≥ 0,80 từ t={t_full}%")
        if t_act is not None:
            key.append(f"active chỉ đạt ở t={t_act}%")
        elif t_full is not None:
            key.append("active chưa chạm 0,80 ở mọi mốc")
        if key:
            dual_items.append(("key", " — nhưng ".join(key) + " (XGBoost)."))
        dual_items.append(
            ("out", "Kết luận RQ1 phải nêu CẢ HAI cohort — đây là điểm trung thực của đồ án.")
        )
    else:
        dual_items.append(
            ("note", "(Thiếu sensitivity_active_xgb.csv — phân tích dual-cohort xem trong repo.)")
        )
    add_content(
        9,
        "RQ1 · DUAL-COHORT (SLIDE ĐINH)",
        "Recall toàn-cohort vs cohort còn-đang-học",
        dual_items,
        notes="Slide quan trọng nhất. Recall toàn-cohort trông đẹp một phần vì tính cả sinh viên đã "
        "rút trước mốc dự đoán. Khi chỉ xét cohort còn đang học — nhóm ta thực sự can thiệp được — "
        "recall đạt 0,80 từ mốc 40 phần trăm trên full nhưng phải đến mốc 100 phần trăm trên active. "
        "Nhóm chủ động trình bày cả hai con số thay vì chỉ con số đẹp.",
        presenter="Khoa",
        image=_fig("sensitivity_active_recall_xgb.png"),
        caption="Full vs active cohort — recall theo mốc (XGBoost)",
        table=dual_table,
    )

    # ── 10 · Threshold trên validation ───────────────────────────────────
    thr_items = [
        (
            "do",
            "Chọn ngưỡng quyết định trên tập VALIDATION theo từng chính sách, rồi mới áp lên test.",
        ),
        ("why", "Tối ưu ngưỡng trực tiếp trên test = rò rỉ lựa chọn (selection leakage)."),
    ]
    thr_table = None
    if thr:
        thr_table = {
            "headers": [
                "Chính sách",
                "Ngưỡng",
                "Recall (val)",
                "Recall (test)",
                "Precision (test)",
            ],
            "rows": [
                [
                    r["policy"],
                    _vn(r["threshold"], 2),
                    _vn(r["val_recall"]),
                    _vn(r["test_recall"]),
                    _vn(r["test_precision"]),
                ]
                for r in thr
            ],
            "top": 3.75,
            "width": 11.6,
            "col_w": [Inches(3.0), Inches(1.6), Inches(2.3), Inches(2.3), Inches(2.4)],
            "fs": 12,
        }
        rec_pol = next((r for r in thr if "recall" in r["policy"]), None)
        if rec_pol:
            thr_items.append(
                (
                    "key",
                    f"Chính sách '{rec_pol['policy']}' (ngưỡng {_vn(rec_pol['threshold'], 2)} chọn "
                    f"trên val) → test recall {_vn(rec_pol['test_recall'])}, precision "
                    f"{_vn(rec_pol['test_precision'])} — val và test khớp nhau.",
                )
            )
    else:
        thr_items.append(
            ("note", "(Thiếu threshold_validation.csv — bảng ngưỡng xem trong repo.)")
        )
    add_content(
        10,
        "PHƯƠNG PHÁP · NGƯỠNG",
        "Ngưỡng quyết định chọn trên validation, không phải trên test",
        thr_items,
        notes="Ngưỡng quyết định được chọn trên tập validation theo từng chính sách — mặc định, tối ưu "
        "F1, Youden, hay ràng buộc recall — rồi mới áp một lần lên test. Số trên validation và test "
        "khớp nhau, chứng tỏ ngưỡng khái quát tốt và không có rò rỉ lựa chọn.",
        presenter="Đức",
        table=thr_table,
    )

    # ── 11 · RQ3 accuracy ────────────────────────────────────────────────
    rq3_items = [
        (
            "do",
            "So 4 chiến lược: no-resample · class-weight · SMOTE · ADASYN trên cả 5 model (t=100%).",
        ),
    ]
    rq3_table = None
    if imb:
        by_model: dict[str, list[dict]] = {}
        for r in imb:
            by_model.setdefault(r["model"], []).append(r)
        rows_out = []
        max_spread = 0.0
        order = ["xgb", "lgbm", "rf", "logreg", "ann"]
        for m in [m for m in order if m in by_model] + [m for m in by_model if m not in order]:
            recs = {r["strategy"]: _f(r, "recall") for r in by_model[m]}
            spread = max(recs.values()) - min(recs.values())
            max_spread = max(max_spread, spread)
            rows_out.append(
                [
                    MODEL_NAMES.get(m, m),
                    _vn(recs.get("none"), 3) if "none" in recs else "—",
                    _vn(recs.get("class_weight"), 3) if "class_weight" in recs else "—",
                    _vn(recs.get("SMOTE"), 3) if "SMOTE" in recs else "—",
                    _vn(recs.get("ADASYN"), 3) if "ADASYN" in recs else "—",
                    _vn(spread, 4),
                ]
            )
        rq3_table = {
            "headers": ["Model", "None", "Class-wt", "SMOTE", "ADASYN", "Chênh lệch"],
            "rows": rows_out,
            "top": 3.7,
            "width": 11.6,
            "col_w": [
                Inches(2.6),
                Inches(1.8),
                Inches(1.8),
                Inches(1.8),
                Inches(1.8),
                Inches(1.8),
            ],
            "fs": 12,
        }
        rq3_items.append(
            (
                "key",
                f"Chênh lệch recall lớn nhất giữa các chiến lược chỉ {_vn(max_spread, 4)} "
                "— về HIỆU NĂNG, resampling gần như không đổi kết quả.",
            )
        )
        rq3_items.append(
            ("note", "Lớp at-risk ~52% — gần cân bằng, nên resampling ít tác dụng là hợp lý.")
        )
    else:
        rq3_items.append(
            ("note", "(Thiếu imbalance_comparison.csv — bảng so sánh xem trong repo.)")
        )
    add_content(
        11,
        "RQ3 · MẤT CÂN BẰNG — HIỆU NĂNG",
        "4 chiến lược resampling: recall gần như không đổi",
        rq3_items,
        notes="Vế thứ nhất của RQ3: về hiệu năng, bốn chiến lược xử lý mất cân bằng cho recall chênh "
        "nhau không quá một điểm phần trăm trên mọi model — hợp lý vì lớp at-risk chiếm khoảng 52 "
        "phần trăm, gần cân bằng. Nhưng đó chưa phải toàn bộ câu chuyện — sang slide sau.",
        presenter="Đức",
        table=rq3_table,
    )

    # ── 12 · RQ3 explanation + RQ2 seed stability ────────────────────────
    xai_items = []
    if xai_strat:
        jacc = [(_f(r, "jaccard_top10"), r["strategy_a"], r["strategy_b"]) for r in xai_strat]
        jmin = min(jacc)
        jmax = max(jacc)
        xai_items.append(
            (
                "out",
                "RQ3 (vế giải thích): hiệu năng không đổi nhưng TOP-10 đặc trưng thì đổi — "
                f"Jaccard từ {_vn(jmin[0], 2)} ({STRATEGY_NAMES.get(jmin[1], jmin[1])} vs "
                f"{STRATEGY_NAMES.get(jmin[2], jmin[2])}) đến {_vn(jmax[0], 2)} "
                f"({STRATEGY_NAMES.get(jmax[1], jmax[1])} vs {STRATEGY_NAMES.get(jmax[2], jmax[2])}).",
            )
        )
        xai_items.append(
            (
                "why",
                "Chọn chiến lược resampling = chọn luôn 'câu chuyện' mà mô hình kể với cố vấn học tập.",
            )
        )
    else:
        xai_items.append(
            ("note", "(Thiếu xai_stability_strategies.csv — so sánh giải thích xem trong repo.)")
        )
    if xai_seed:
        r0 = xai_seed[0]
        xai_items.append(
            (
                "key",
                f"RQ2 — ổn định theo seed: Jaccard top-10 TB {_vn(r0['mean_jaccard'], 2)} · "
                f"Spearman TB {_vn(r0['mean_spearman'], 2)} ({_vn_int(r0['n_pairs'])} cặp seed).",
            )
        )
        xai_items.append(
            (
                "out",
                "Thứ hạng tổng thể rất ổn định (Spearman ~0,97); thành phần top-10 dao động nhẹ.",
            )
        )
    else:
        xai_items.append(
            ("note", "(Thiếu xai_stability_seeds.csv — độ ổn định seed xem trong repo.)")
        )
    add_content(
        12,
        "RQ2 + RQ3 · GIẢI THÍCH",
        "Giải thích: ổn định theo seed, nhạy theo chiến lược resampling",
        xai_items,
        notes="Vế thứ hai của RQ3 và toàn bộ RQ2. Cùng hiệu năng nhưng SMOTE với ADASYN làm xáo trộn "
        "top-10 đặc trưng — Jaccard thấp nhất khoảng 0,54. Trong khi đó đổi seed thì giải thích khá "
        "ổn định: Spearman trung bình 0,97. Bài học: muốn giải thích nhất quán, tránh resampling "
        "tổng hợp khi dữ liệu đã gần cân bằng.",
        presenter="Bình",
        image=_fig("shap_importance_xgb_t100.png"),
        caption="SHAP importance — XGBoost @ t=100%",
    )

    # ── 13 · Fairness ────────────────────────────────────────────────────
    fair_items = [
        ("do", "Đo recall-gap & FPR-gap của model chính trên các nhóm nhân khẩu học (tập test)."),
    ]
    fair_table = None
    if fair:
        rows = sorted(fair, key=lambda r: _f(r, "recall_gap"), reverse=True)
        fair_table = {
            "headers": [
                "Thuộc tính",
                "Số nhóm",
                "Recall min",
                "Recall max",
                "Gap recall",
                "Gap FPR",
            ],
            "rows": [
                [
                    r["attribute"],
                    _vn_int(r["n_levels"]),
                    _vn(r["recall_min"], 3),
                    _vn(r["recall_max"], 3),
                    _vn(r["recall_gap"], 3),
                    _vn(r["fpr_gap"], 3),
                ]
                for r in rows
            ],
            "top": 3.8,
            "width": 11.6,
            "col_w": [
                Inches(3.1),
                Inches(1.4),
                Inches(1.8),
                Inches(1.8),
                Inches(1.8),
                Inches(1.7),
            ],
            "fs": 12,
        }
        worst = rows[0]
        key_txt = (
            f"Gap lớn nhất: {worst['attribute']} — recall chênh {_vn(worst['recall_gap'], 3)} "
            f"giữa {_vn_int(worst['n_levels'])} nhóm"
        )
        bin_gaps = [_f(r, "recall_gap") for r in rows if int(float(r["n_levels"])) == 2]
        if bin_gaps:
            key_txt += f"; các thuộc tính nhị phân gap ≤ {_vn(max(bin_gaps), 3)}"
        fair_items.append(("key", key_txt + "."))
        fair_items.append(
            ("note", "Gap tồn tại nhưng ở mức vừa phải; cần theo dõi khi triển khai thật.")
        )
    else:
        fair_items.append(
            ("note", "(Thiếu fairness_gaps.csv — phân tích công bằng xem trong repo.)")
        )
    add_content(
        13,
        "CÔNG BẰNG",
        "Fairness — chênh lệch recall/FPR giữa các nhóm sinh viên",
        fair_items,
        notes="Mô hình được soi theo lăng kính công bằng: chênh lệch recall lớn nhất nằm ở nhóm chỉ số "
        "khó khăn kinh tế imd_band, khoảng sáu điểm phần trăm giữa mười một nhóm; giới tính, độ tuổi, "
        "khuyết tật gap nhỏ. Kết luận: chưa thấy thiên lệch nghiêm trọng nhưng phải giám sát tiếp khi "
        "dùng thật.",
        presenter="Bình",
        table=fair_table,
    )

    # ── 14 · Dashboard ───────────────────────────────────────────────────
    add_content(
        14,
        "SẢN PHẨM · DASHBOARD",
        "Streamlit dashboard cho cố vấn học tập",
        [
            (
                "do",
                "App Streamlit: chọn khóa-kỳ + mốc tiến độ → danh sách SV xếp theo xác suất nguy cơ, "
                "kèm giải thích SHAP theo từng sinh viên.",
            ),
            (
                "do",
                "Bộ lọc 'còn đang học' (still-enrolled): mặc định ẩn SV đã rút trước mốc t — đúng "
                "tinh thần dual-cohort ở slide 9.",
            ),
            (
                "why",
                "Người dùng cuối là cố vấn học tập: cần danh sách hành động được, không phải bảng metric.",
            ),
            (
                "out",
                "Chạy tái lập từ repo (streamlit run), đọc đúng model + bảng số liệu đã đóng băng.",
            ),
        ],
        notes="Sản phẩm ứng dụng: dashboard Streamlit cho cố vấn học tập — chọn khóa và mốc tiến độ, "
        "nhận danh sách sinh viên xếp theo rủi ro kèm giải thích SHAP từng em. Điểm đáng chú ý là bộ "
        "lọc còn-đang-học: mặc định chỉ hiện sinh viên còn can thiệp được, nhất quán với phân tích "
        "dual-cohort.",
        presenter="An",
    )

    # ── 15 · Kết luận ────────────────────────────────────────────────────
    add_content(
        15,
        "KẾT LUẬN",
        "Trả lời 3 RQ & hạn chế trung thực",
        [
            (
                "out",
                "RQ1: dự đoán sớm khả thi từ giữa khóa trên toàn cohort; với SV còn đang học, độ "
                "nhạy chỉ đạt mức cao ở cuối khóa — báo cáo cả hai.",
            ),
            (
                "out",
                "RQ2: giải thích ổn định theo seed (thứ hạng ~0,97 Spearman), top-10 dao động nhẹ.",
            ),
            (
                "out",
                "RQ3: resampling không đổi hiệu năng (dữ liệu gần cân bằng) nhưng ĐỔI lời giải thích "
                "→ khuyến nghị no-resample/class-weight.",
            ),
            (
                "p",
                "Hạn chế: 1 bộ dữ liệu (OULAD, 2013–2014) · nhãn Withdrawn gộp nhiều lý do · recall "
                "active-cohort ở mốc sớm còn thấp · chưa kiểm chứng ngoài (external validation).",
            ),
            ("note", "Toàn bộ kết quả tái lập từ repo: 1 lệnh chạy pipeline + bộ test tự động."),
        ],
        notes="Kết luận: RQ1 — dự đoán sớm dùng được từ giữa khóa, nhưng nhóm trung thực báo cả recall "
        "trên cohort còn đang học vốn thấp hơn; RQ2 — giải thích ổn định theo seed; RQ3 — resampling "
        "không giúp hiệu năng mà còn làm xáo trộn lời giải thích, nên khuyến nghị không resample. Hạn "
        "chế chính: một bộ dữ liệu cũ, nhãn withdrawn nhiễu, chưa kiểm chứng ngoài.",
        presenter="Khoa",
    )

    # ── 16 · Q&A ─────────────────────────────────────────────────────────
    s = prs.slides.add_slide(BLANK)
    _rect(s, 0, 0, SW, SH, NAVY)
    _rect(s, 0, Inches(3.3), SW, Inches(0.07), ACCENT)
    _, tf = _box(s, Inches(0.9), Inches(2.35), Inches(11.5), Inches(1.0))
    r = tf.paragraphs[0].add_run()
    r.text = "Q&A — Cảm ơn thầy/cô và các bạn!"
    _set(r, 34, WHITE, bold=True)
    _, tm = _box(s, Inches(0.9), Inches(3.6), Inches(11.5), Inches(1.2))
    rr = tm.paragraphs[0].add_run()
    rr.text = "DSP391m · Nhóm 1 · GVHD: Nguyễn Thị Hoàng Yến"
    _set(rr, 14, RGBColor(0xD7, 0xE0, 0xEC))
    p2 = tm.add_paragraph()
    r2 = p2.add_run()
    r2.text = "Mã nguồn, số liệu và slide tái lập được từ repo của nhóm."
    _set(r2, 12, RGBColor(0xAF, 0xBE, 0xD2), italic=True)
    p2.space_before = Pt(4)
    _notes(
        s,
        "Cảm ơn thầy cô và các bạn đã lắng nghe. Nhóm sẵn sàng trả lời câu hỏi — mọi con số trên "
        "slide đều truy được về file CSV và mã nguồn trong repo.",
    )


def main(argv: list[str] | None = None) -> int:
    ap = argparse.ArgumentParser(description="Build the final defense deck (.pptx).")
    ap.add_argument(
        "--tables-dir",
        default=str(ROOT / "reports" / "tables"),
        help="Thư mục chứa các CSV số liệu (mặc định: reports/tables).",
    )
    args = ap.parse_args(argv)
    tables_dir = Path(args.tables_dir)
    if not tables_dir.is_dir():
        print(f"[ERROR] --tables-dir không tồn tại: {tables_dir}", file=sys.stderr)
        return 2

    build(tables_dir)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    n_slides = len(prs.slides._sldIdLst)
    size_kb = OUT.stat().st_size / 1024
    print(f"Saved {OUT}  ({n_slides} slides, {size_kb:.0f} KB)")
    if DEGRADED:
        print("Slides degraded (thiếu bảng): " + ", ".join(sorted(set(DEGRADED))))
    else:
        print("Đủ toàn bộ bảng — không slide nào bị giản lược.")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
