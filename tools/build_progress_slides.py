"""Build the current-progress presentation (.pptx) for the DSP391m project.

Generates ``reports/slides/Progress_Report.pptx`` (16:9) — a status deck covering
the whole pipeline as it stands: data (Phase 1), modeling + imbalance (Phase 2),
time-aware RQ1 (Phase 3), fine-tuning (Task 4), the supplementary evaluation
(confusion matrix / threshold tuning / significance tests), and explainability
(Phase 5: SHAP/LIME + stability RQ2/RQ3) — plus an honest what's-left slide.

Reuses the house style of ``build_task3_slides.py`` and embeds the real figures
from ``reports/figures``. Speaker notes go into each slide.

    python tools/build_progress_slides.py
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

try:
    from PIL import Image

    def _img_size(path: Path) -> tuple[int, int]:
        with Image.open(path) as im:
            return im.size
except Exception:  # pragma: no cover - PIL optional

    def _img_size(path: Path) -> tuple[int, int]:
        return (1600, 900)


ROOT = Path(__file__).resolve().parents[1]
FIG = ROOT / "reports" / "figures"
OUT = ROOT / "reports" / "slides" / "Progress_Report.pptx"

# ── Palette (same as the Task 3 deck) ───────────────────────────────────────
NAVY = RGBColor(0x0F, 0x2A, 0x4F)
ACCENT = RGBColor(0xE8, 0x72, 0x2B)
GREEN = RGBColor(0x1E, 0x8E, 0x4E)
BLUE = RGBColor(0x1F, 0x6F, 0xB2)
AMBER = RGBColor(0xB9, 0x7A, 0x0C)
DARK = RGBColor(0x26, 0x2B, 0x33)
GREY = RGBColor(0x5B, 0x63, 0x70)
LIGHT = RGBColor(0xF2, 0xF4, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Calibri"
SW, SH = Inches(13.333), Inches(7.5)
LABELS = {"do": ("Làm gì", GREEN), "why": ("Vì sao", BLUE), "out": ("Kết quả", AMBER)}

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]


def _set(run, size, color, bold=False, italic=False, font=FONT):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font


def _box(slide, left, top, width, height):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    return tb, tf


def _rect(slide, left, top, width, height, fill, line=None):
    from pptx.enum.shapes import MSO_SHAPE

    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


def _footer(slide, n, presenter=""):
    _rect(slide, 0, 0, SW, Inches(0.16), NAVY)
    _rect(slide, 0, 0, Inches(3.2), Inches(0.16), ACCENT)
    _, tf = _box(slide, Inches(0.55), Inches(7.06), Inches(9.5), Inches(0.34))
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "DSP391m · Nhóm 1 · Báo cáo tiến độ dự án"
    _set(r, 9, GREY)
    if presenter:
        r2 = p.add_run()
        r2.text = f"   ·   {presenter}"
        _set(r2, 9, ACCENT, bold=True)
    _, tf2 = _box(slide, Inches(12.2), Inches(7.06), Inches(0.9), Inches(0.34))
    tf2.paragraphs[0].alignment = PP_ALIGN.RIGHT
    rr = tf2.paragraphs[0].add_run()
    rr.text = str(n)
    _set(rr, 10, NAVY, bold=True)


def _notes(slide, text):
    if text:
        slide.notes_slide.notes_text_frame.text = text


def add_title(title, subtitle, meta, presenters):
    s = prs.slides.add_slide(BLANK)
    _rect(s, 0, 0, SW, SH, NAVY)
    _rect(s, 0, Inches(2.55), SW, Inches(0.07), ACCENT)
    _, tf = _box(s, Inches(0.9), Inches(2.7), Inches(11.5), Inches(2.2))
    r = tf.paragraphs[0].add_run()
    r.text = title
    _set(r, 33, WHITE, bold=True)
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = subtitle
    _set(r2, 16, RGBColor(0xCF, 0xDC, 0xEC), italic=True)
    p2.space_before = Pt(10)
    _, tk = _box(s, Inches(0.9), Inches(1.5), Inches(11.0), Inches(0.7))
    rk = tk.paragraphs[0].add_run()
    rk.text = "BÁO CÁO TIẾN ĐỘ"
    _set(rk, 15, ACCENT, bold=True)
    _, tm = _box(s, Inches(0.9), Inches(5.3), Inches(11.5), Inches(1.7))
    for i, line in enumerate(meta):
        pp = tm.paragraphs[0] if i == 0 else tm.add_paragraph()
        rr = pp.add_run()
        rr.text = line
        _set(rr, 13, RGBColor(0xD7, 0xE0, 0xEC))
        pp.space_after = Pt(2)
    pp = tm.add_paragraph()
    rr = pp.add_run()
    rr.text = presenters
    _set(rr, 12.5, WHITE, bold=True)
    pp.space_before = Pt(6)


def add_divider(part, title, subtitle, n):
    s = prs.slides.add_slide(BLANK)
    _rect(s, 0, 0, SW, SH, NAVY)
    _rect(s, Inches(0.9), Inches(2.7), Inches(1.6), Inches(0.09), ACCENT)
    _, tk = _box(s, Inches(0.9), Inches(2.0), Inches(11), Inches(0.7))
    rk = tk.paragraphs[0].add_run()
    rk.text = part
    _set(rk, 15, ACCENT, bold=True)
    _, tf = _box(s, Inches(0.9), Inches(2.95), Inches(11.3), Inches(2.0))
    r = tf.paragraphs[0].add_run()
    r.text = title
    _set(r, 30, WHITE, bold=True)
    if subtitle:
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = subtitle
        _set(r2, 15, RGBColor(0xCF, 0xDC, 0xEC), italic=True)
        p2.space_before = Pt(8)


def _add_heading(s, kicker, title):
    _, tk = _box(s, Inches(0.55), Inches(0.32), Inches(12), Inches(0.4))
    rk = tk.paragraphs[0].add_run()
    rk.text = kicker
    _set(rk, 12, ACCENT, bold=True)
    _, tt = _box(s, Inches(0.55), Inches(0.66), Inches(12.2), Inches(0.95))
    rt = tt.paragraphs[0].add_run()
    rt.text = title
    _set(rt, 25, NAVY, bold=True)
    _rect(s, Inches(0.55), Inches(1.62), Inches(12.2), Inches(0.025), RGBColor(0xD8, 0xDF, 0xE8))


def _render_bullets(tf, items):
    first = True
    for it in items:
        kind = it[0]
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(7)
        p.line_spacing = 1.04
        if kind in LABELS:
            label, col = LABELS[kind]
            dot = p.add_run()
            dot.text = "● "
            _set(dot, 13, col, bold=True)
            lab = p.add_run()
            lab.text = f"{label}: "
            _set(lab, 14.5, col, bold=True)
            txt = p.add_run()
            txt.text = it[1]
            _set(txt, 14.5, DARK)
        elif kind == "p":
            dot = p.add_run()
            dot.text = "▪ "
            _set(dot, 12, NAVY, bold=True)
            txt = p.add_run()
            txt.text = it[1]
            _set(txt, 14.5, DARK)
        elif kind == "s":
            p.level = 1
            dot = p.add_run()
            dot.text = "– "
            _set(dot, 12, GREY)
            txt = p.add_run()
            txt.text = it[1]
            _set(txt, 12.5, GREY)
            p.space_after = Pt(4)
        elif kind == "key":
            bar = p.add_run()
            bar.text = "▶ Số liệu chốt:  "
            _set(bar, 13.5, ACCENT, bold=True)
            txt = p.add_run()
            txt.text = it[1]
            _set(txt, 13.5, NAVY, bold=True)
            p.space_before = Pt(6)
        elif kind == "note":
            txt = p.add_run()
            txt.text = it[1]
            _set(txt, 12, GREY, italic=True)
            p.space_before = Pt(4)


def _add_image(s, path: Path, left, top, max_w, max_h, caption=""):
    w, h = _img_size(path)
    ratio = min(max_w / w, max_h / h)
    iw, ih = int(w * ratio), int(h * ratio)
    pic_left = left + (max_w - iw) // 2
    s.shapes.add_picture(str(path), pic_left, top, width=Emu(iw), height=Emu(ih))
    if caption:
        _, tf = _box(s, left, top + Emu(ih) + Inches(0.05), Emu(max_w), Inches(0.4))
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        r = tf.paragraphs[0].add_run()
        r.text = caption
        _set(r, 9.5, GREY, italic=True)


def _add_table(s, headers, rows, left, top, width, col_w=None, fontsize=11):
    nrows, ncols = len(rows) + 1, len(headers)
    height = Inches(0.36 * nrows)
    gtbl = s.shapes.add_table(nrows, ncols, left, top, width, height)
    tbl = gtbl.table
    if col_w:
        for i, cw in enumerate(col_w):
            tbl.columns[i].width = cw
    for j, htext in enumerate(headers):
        c = tbl.cell(0, j)
        c.fill.solid()
        c.fill.fore_color.rgb = NAVY
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        c.margin_top = c.margin_bottom = Pt(2)
        r = c.text_frame.paragraphs[0].add_run()
        r.text = htext
        _set(r, fontsize, WHITE, bold=True)
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            c = tbl.cell(i, j)
            c.fill.solid()
            c.fill.fore_color.rgb = WHITE if i % 2 else LIGHT
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            c.margin_top = c.margin_bottom = Pt(1)
            r = c.text_frame.paragraphs[0].add_run()
            r.text = str(val)
            _set(r, fontsize, NAVY if j == 0 else DARK, bold=(j == 0))
    return gtbl


def add_content(
    n,
    kicker,
    title,
    items,
    notes="",
    presenter="",
    image=None,
    caption="",
    table=None,
    text_w=12.2,
):
    s = prs.slides.add_slide(BLANK)
    _footer(s, n, presenter)
    _add_heading(s, kicker, title)
    if image:
        text_w = 6.85
    _, tf = _box(s, Inches(0.6), Inches(1.85), Inches(text_w), Inches(5.0))
    tf.vertical_anchor = MSO_ANCHOR.TOP
    _render_bullets(tf, items)
    if image:
        _add_image(s, image, Inches(7.7), Inches(1.95), Inches(5.15), Inches(4.6), caption)
    if table:
        _add_table(
            s,
            table["headers"],
            table["rows"],
            Inches(table.get("left", 0.6)),
            Inches(table["top"]),
            Inches(table.get("width", 9.5)),
            col_w=table.get("col_w"),
            fontsize=table.get("fs", 11),
        )
    _notes(s, notes)
    return s


# ════════════════════════════════════════════════════════════════════════
# SLIDES
# ════════════════════════════════════════════════════════════════════════

add_title(
    "Time-Aware Explainable ML — phát hiện sớm sinh viên nguy cơ (OULAD)",
    "Báo cáo tiến độ: dữ liệu → mô hình hoá → giải thích (Phase 1–5)",
    [
        "DSP391m · Đồ án Khoa học Dữ liệu · Nhóm 1 · Đại học FPT",
        "Cập nhật: 02/07/2026",
    ],
    "Thành viên: Sơn · Khoa · An · Đức · Phúc · Bình",
)

add_content(
    2,
    "TỔNG QUAN",
    "Tiến độ theo pha & người phụ trách",
    [
        ("note", "✅ đã có code + kết quả thật   ·   🔄 đang làm   ·   ⏳ chưa bắt đầu"),
    ],
    notes="Đây là bức tranh tiến độ tổng thể. Toàn bộ khối dữ liệu và mô hình hoá — Phase 1 đến "
    "Phase 5 — đã có code và kết quả thật. Còn lại là dashboard của An, phần literature review của "
    "Sơn, và bản báo cáo/slide cuối đang được tổng hợp.",
    table={
        "headers": ["Pha / Nhiệm vụ", "Người", "Trạng thái"],
        "rows": [
            ["Phase 1 — Data pipeline, cắt mốc, harness", "Phúc", "✅ Hoàn thành"],
            ["Phase 2 — Benchmark 5 model + SMOTE", "Đức", "✅ Có kết quả"],
            ["Phase 3 — Time-aware RQ1 (6 mốc)", "Khoa", "✅ Có kết quả"],
            ["Task 4 — Fine-tuning + sensitivity", "Đức/Khoa", "✅ Hoàn thành"],
            ["Bổ trợ — Confusion / Threshold / Kiểm định", "Nhóm", "✅ Mới bổ sung"],
            ["Phase 5 — SHAP/LIME + stability (RQ2/RQ3)", "Bình", "✅ Có kết quả"],
            ["Phase 6a — Dashboard Streamlit", "An", "⏳ Chưa làm"],
            ["Literature Review / References", "Sơn", "🔄 Tài liệu nền"],
            ["Báo cáo & slide cuối", "Nhóm", "🔄 Đang tổng hợp"],
        ],
        "top": 2.05,
        "width": 12.2,
        "col_w": [Inches(7.2), Inches(2.0), Inches(3.0)],
        "fs": 12.5,
    },
)

add_content(
    3,
    "PHASE 1 · DỮ LIỆU",
    "Pipeline dữ liệu chống rò rỉ (nền tảng)",
    [
        (
            "do",
            "7 bảng OULAD → master (32.593 × 33) → cắt 6 mốc 10–100% → tiền xử lý fit-on-train.",
        ),
        ("why", "Mọi bước mô hình hoá đứng trên tập test cố định + pipeline chống rò rỉ này."),
        ("out", "6 bộ dataset_t10…t100 + test 20% cố định; 19/19 kiểm thử rò rỉ ĐẠT."),
        ("key", "Test cố định dùng lại y hệt qua 6 mốc → 6 điểm hiệu năng SO SÁNH ĐƯỢC."),
    ],
    presenter="Phúc",
    notes="Nền tảng là pipeline dữ liệu của Phúc: gộp 7 bảng thành master 32.593 dòng, cắt theo 6 mốc "
    "tiến độ, tiền xử lý fit-on-train chống rò rỉ, và một tập test cố định. 19 kiểm thử tự động xác "
    "nhận không rò rỉ. Đây là điều kiện để 6 điểm hiệu năng theo mốc so sánh được.",
    image=(FIG / "preprocessing_sequence.png"),
    caption="Trình tự pipeline tiền xử lý",
)

add_content(
    4,
    "PHASE 2 · MÔ HÌNH",
    "Benchmark 5 thuật toán + xử lý mất cân bằng",
    [
        ("do", "LR · RF · XGBoost · LightGBM · ANN; SMOTE cân bằng CHỈ trên train fold."),
        ("why", "Bỏ sót sinh viên nguy cơ là sai lầm đắt nhất → recall & PR-AUC là chỉ số chính."),
        ("out", "Bảng model_metrics.csv (30 dòng: 5 model × 6 mốc) + biểu đồ so sánh."),
        ("key", "Tại t=100%: XGB recall 0,935 · PR-AUC 0,991; các model tree bám sát nhau."),
    ],
    presenter="Đức",
    notes="Phase 2 benchmark năm thuật toán với SMOTE cân bằng chỉ trên tập train để tránh rò rỉ. Vì bỏ "
    "sót sinh viên nguy cơ là tốn kém nhất, chỉ số chính là recall và PR-AUC của lớp at-risk. Tại "
    "cuối khóa, XGBoost đạt recall 0,935 và PR-AUC 0,991.",
    image=(FIG / "model_benchmark.png"),
    caption="So sánh model tại t=100% (recall/F1/PR-AUC/ROC-AUC)",
)

add_content(
    5,
    "PHASE 3 · RQ1",
    "Hiệu năng theo thời gian — dự đoán được sớm cỡ nào?",
    [
        ("do", "Huấn luyện lại tại 6 mốc; theo dõi recall/PR-AUC theo tiến độ khóa học."),
        ("why", "RQ1: cần cảnh báo SỚM mới can thiệp kịp, không đợi hết môn."),
        (
            "out",
            "Hiệu năng tăng đều theo mốc; đạt ngưỡng tin cậy (recall≥0,80 & PR-AUC≥0,90) từ t=40%.",
        ),
        ("key", "Mốc tin cậy sớm nhất: 40% tiến độ khóa học."),
    ],
    presenter="Khoa",
    notes="Phase 3 trả lời RQ1: dự đoán sớm được đến đâu. Huấn luyện lại ở cả 6 mốc, hiệu năng tăng đều "
    "theo tiến độ. Ngưỡng tin cậy recall trên 0,80 và PR-AUC trên 0,90 đạt được ngay từ mốc 40% "
    "tiến độ khóa học — nghĩa là có thể cảnh báo hữu ích từ khá sớm.",
    image=(FIG / "time_aware_recall.png"),
    caption="Recall theo 6 mốc thời gian (RQ1)",
)

add_content(
    6,
    "BỔ TRỢ · KIỂM ĐỊNH",
    "Xếp hạng model có ý nghĩa thống kê không?",
    [
        (
            "do",
            "Friedman trên 25 fold CV ghép cặp (mỗi metric) + post-hoc Wilcoxon (hiệu chỉnh Holm).",
        ),
        ("why", "Bảng held-out chỉ có 1 số/model → không phân biệt được thắng thật với may rủi."),
        (
            "out",
            "Mọi metric p < 1e-12 → khác biệt CÓ ý nghĩa; XGB nhất recall, LGBM nhất F1/PR-AUC.",
        ),
        ("key", "XGB thắng recall có ý nghĩa so với cả 4 model còn lại (Holm p<0,05)."),
    ],
    presenter="Nhóm",
    notes="Phần bổ trợ đầu tiên: kiểm định thống kê. Dùng Friedman trên 25 fold cross-validation ghép "
    "cặp cho từng metric, rồi post-hoc Wilcoxon hiệu chỉnh Holm. Kết quả mọi metric p rất nhỏ, dưới "
    "1e-12: khác biệt là thật. XGBoost đứng nhất về recall và thắng có ý nghĩa so với cả bốn model "
    "còn lại; LightGBM nhất về F1 và PR-AUC.",
    table={
        "headers": ["Metric", "Model tốt nhất", "Friedman p"],
        "rows": [
            ["recall", "XGBoost", "2,4e-13"],
            ["F1", "LightGBM", "6,1e-17"],
            ["PR-AUC", "LightGBM", "3,9e-17"],
            ["ROC-AUC", "LightGBM", "1,4e-16"],
        ],
        "top": 4.7,
        "width": 6.7,
        "left": 0.6,
        "col_w": [Inches(2.2), Inches(2.6), Inches(1.9)],
        "fs": 12,
    },
)

add_content(
    7,
    "BỔ TRỢ · NGƯỠNG",
    "Tinh chỉnh ngưỡng quyết định (threshold tuning)",
    [
        ("do", "Quét ngưỡng; chọn theo max-F1, Youden's J, và chính sách recall≥0,90."),
        ("why", "Ngưỡng mặc định 0,5 là tuỳ tiện cho bài toán mất cân bằng/cảnh báo sớm."),
        (
            "out",
            "Chính sách recall≥0,90 → ngưỡng 0,87 nâng precision lên 0,995 mà vẫn giữ recall 0,90.",
        ),
        ("key", "0,5 → 0,87: precision 0,973 → 0,995 (ít báo động giả hơn nhiều)."),
    ],
    presenter="Nhóm",
    notes="Phần bổ trợ thứ hai: tinh chỉnh ngưỡng. Ngưỡng mặc định 0,5 là tuỳ tiện. Nhóm quét ngưỡng và "
    "chọn theo ba chính sách. Với chính sách giữ recall tối thiểu 0,90, ngưỡng tối ưu là 0,87, nâng "
    "precision từ 0,973 lên 0,995 — giảm mạnh báo động giả — mà vẫn bắt được 90% sinh viên nguy cơ.",
    image=(FIG / "threshold_tuning.png"),
    caption="Precision/Recall/F1 theo ngưỡng",
)

add_content(
    8,
    "BỔ TRỢ · CONFUSION",
    "Ma trận nhầm lẫn tại ngưỡng đã tinh chỉnh",
    [
        ("do", "Confusion matrix của XGB @ t=100% ở ngưỡng tuned 0,87."),
        ("why", "Một con số recall che giấu ‘lỗi nằm ở đâu’ — ma trận cho thấy FN/FP cụ thể."),
        ("out", "Bỏ sót (FN) 10% sinh viên nguy cơ; báo động giả (FP) chỉ 0,4%."),
        ("key", "FN 337/3.376 · FP 14/3.113 — cân bằng phù hợp cho cảnh báo sớm."),
    ],
    presenter="Nhóm",
    notes="Ma trận nhầm lẫn cho thấy lỗi nằm ở đâu. Ở ngưỡng tinh chỉnh 0,87, mô hình bỏ sót 10% sinh "
    "viên nguy cơ và chỉ báo động giả 0,4%. Đây là cân bằng phù hợp khi năng lực can thiệp có hạn.",
    image=(FIG / "confusion_tuned_t100.png"),
    caption="XGB @ t=100%, ngưỡng 0,87",
)

add_content(
    9,
    "PHASE 5 · XAI",
    "Giải thích mô hình — SHAP",
    [
        ("do", "SHAP TreeExplainer: tầm quan trọng toàn cục + hướng tác động của từng đặc trưng."),
        ("why", "Recall 0,93 vô nghĩa nếu giáo viên không biết VÌ SAO sinh viên bị gắn cờ."),
        (
            "out",
            "Top driver: số ngày không hoạt động, điểm tích luỹ (có/không trọng số), số bài đã nộp.",
        ),
        ("key", "days_since_last_activity cao → đẩy mạnh về ‘at-risk’ (khớp trực giác sư phạm)."),
    ],
    presenter="Bình",
    notes="Phase 5 mở hộp đen. SHAP cho thấy các yếu tố chi phối dự đoán: số ngày kể từ lần hoạt động "
    "cuối, điểm tích luỹ, và số bài đã nộp. Biểu đồ beeswarm cho thấy nhiều ngày không hoạt động đẩy "
    "dự đoán về phía nguy cơ — khớp trực giác sư phạm.",
    image=(FIG / "shap_summary_xgb_t100.png"),
    caption="SHAP summary (beeswarm) — XGB @ t=100%",
)

add_content(
    10,
    "PHASE 5 · RQ2/RQ3",
    "Độ ổn định của giải thích — đóng góp nghiên cứu",
    [
        ("do", "LIME giải thích từng sinh viên; đo ổn định theo SEED (RQ2) & theo MỐC (RQ3)."),
        (
            "why",
            "Giải thích chỉ đáng tin nếu không đổi khi seed đổi, và drift có kiểm soát theo thời gian.",
        ),
        (
            "out",
            "RQ2: rất ổn định (Spearman 0,97 · Jaccard 0,75). RQ3: mốc liền kề bám sát, mốc sớm lệch xa t=100%.",
        ),
        ("key", "SHAP vs LIME khớp ở top driver, lệch ở giữa (Jaccard@10 = 0,25)."),
    ],
    presenter="Bình",
    notes="Đóng góp nghiên cứu là đo độ ổn định của giải thích. RQ2 — đổi seed: rất ổn định, Spearman "
    "0,97. RQ3 — theo mốc thời gian: các mốc liền kề bám sát nhau nhưng mốc sớm lệch xa so với cuối "
    "khóa, nên khi trình bày cho giáo viên phải ghi rõ mốc. SHAP và LIME khớp ở nhóm driver hàng "
    "đầu nhưng lệch ở tầm giữa — lý do nhóm báo cáo cả hai và kiểm định ổn định thay vì tin một cái.",
    image=(FIG / "xai_stability_drift.png"),
    caption="Drift giải thích theo mốc (RQ3)",
)

add_content(
    11,
    "CÒN LẠI",
    "Việc chưa làm & kế hoạch",
    [
        ("p", "An — Dashboard Streamlit + đóng gói model (Phase 6a): CHƯA bắt đầu."),
        ("p", "Sơn — Introduction / Literature Review / References: mới có tài liệu nền."),
        ("p", "Nhóm — Báo cáo cuối + slide toàn dự án: đang tổng hợp (2 notebook 05/06 đã sẵn)."),
        ("note", "Khối dữ liệu + mô hình hoá + XAI đã xong; phần còn lại là đóng gói & viết."),
    ],
    presenter="Nhóm",
    notes="Về việc còn lại: dashboard của An chưa bắt đầu; phần literature review của Sơn mới có tài liệu "
    "nền; và bản báo cáo cùng slide cuối đang được tổng hợp, dựa trên hai notebook mô hình hoá và "
    "giải thích đã hoàn tất. Toàn bộ khối kỹ thuật đã xong, phần còn lại chủ yếu là đóng gói và viết.",
    table={
        "headers": ["Hạng mục", "Người", "Kế hoạch"],
        "rows": [
            ["Dashboard Streamlit", "An", "Dựng app dự đoán + tải model"],
            ["Literature Review", "Sơn", "Viết theo template báo cáo"],
            ["Báo cáo & slide cuối", "Nhóm", "Ghép từ notebook 05/06 + kết quả"],
        ],
        "top": 4.6,
        "width": 12.2,
        "col_w": [Inches(4.6), Inches(1.6), Inches(6.0)],
        "fs": 12,
    },
)

s12 = add_content(
    12,
    "KẾT LUẬN",
    "Tình trạng hiện tại",
    [
        (
            "out",
            "Đã có: pipeline dữ liệu chống rò rỉ · benchmark 5 model · RQ1 theo 6 mốc · fine-tuning.",
        ),
        (
            "out",
            "Bổ sung: confusion matrix · threshold tuning · kiểm định thống kê · SHAP/LIME + stability.",
        ),
        ("do", "Tiếp theo: dashboard (An) · literature review (Sơn) · báo cáo/slide cuối (nhóm)."),
        ("note", "Mọi kết quả tái lập được từ src/ + tools/; 19/19 test đạt."),
    ],
    presenter="Nhóm",
    notes="Tóm lại: khối kỹ thuật đã hoàn chỉnh và tái lập được — pipeline dữ liệu, benchmark, RQ1 theo "
    "mốc, fine-tuning, cùng các phân tích bổ trợ và toàn bộ lớp giải thích SHAP/LIME với đo ổn định. "
    "Bước tiếp theo là dashboard, literature review, và báo cáo cuối. Cảm ơn thầy/cô và cả lớp.",
)
_, tf = _box(s12, Inches(0.6), Inches(6.1), Inches(12.2), Inches(0.8))
r = tf.paragraphs[0].add_run()
r.text = "Cảm ơn đã lắng nghe!"
_set(r, 22, ACCENT, bold=True)


OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(OUT))
print(f"Saved {OUT}  ({len(prs.slides._sldIdLst)} slides)")
